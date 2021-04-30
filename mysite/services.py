import zipfile
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import pandas as pd
import numpy as np
import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib.patches import Patch
from matplotlib import ticker
from IPython.display import display, HTML
import re # quickly find float numbers in a string.
import collections
import os

def data_processing(file):
	filename = file.name.split('.')[0]
	sorted_data = pd.read_excel(file, header = 0)
	sorted_data = sorted_data.fillna(0)
	#display(HTML(sorted_data.to_html()))
	print(sorted_data.columns)
	sorted_data = sorted_data.rename(columns={'Glycans_x000D_\nPos.' : 'Glycans Pos.', 'Sequence\r\n(unformatted)': 'Sequence', 'Calc._x000D_\nMH' : 'Calc. MH', 'PEP_x000D_\n2D':'PEP 2D'})
	#sorted_data.columns

	# print('--------------------------------------- Score & PEP2D Selection Process Begins ---------------------------------------')
	# print('\nStep1: Filter Score & PEP2D -> Delete the peptides without the correct N-glycosylation sequon.\n')

	###########################
	sorted_data_scoreHightoLow_score200_pep2d0001 = sorted_data.loc[((sorted_data['Score'] > 200) & (sorted_data['PEP 2D'] < 0.001))]
	#display(HTML(sorted_data_scoreHightoLow_score200_pep2d0001.to_html()))
	sorted_data_scoreHightoLow_score200_pep2d0001 = sorted_data_scoreHightoLow_score200_pep2d0001.reset_index(drop=True)
	sequence = sorted_data_scoreHightoLow_score200_pep2d0001['Sequence'].tolist()
	all_aa = ['A','R', 'N', 'D', 'B', 'C', 'E', 'Q', 'Z', 'G', 'H', 'I', 'L', 'K', 'M', 'F', 'S', 'T', 'W', 'Y', 'V', 'P']
	pure_seq = []
	for ele in sequence:
	    each_pure_seq = [ word for word in ele if word in all_aa]
	    each_pure_seq = ''.join(each_pure_seq)
	    pure_seq.append(each_pure_seq)
	sorted_data_scoreHightoLow_score200_pep2d0001.insert(sorted_data_scoreHightoLow_score200_pep2d0001.columns.get_loc('Sequence') + 1 , 'Pure Sequence', pure_seq , True)
	# print('The original data size:')
	# print(sorted_data_scoreHightoLow_score200_pep2d0001.shape)
	#display(HTML(sorted_data_scoreHightoLow_score200_pep2d0001.to_html()))
	#sorted_data_scoreHightoLow_score200_pep2d0001.to_excel('%s_sorted_data_score200_pep2d0001.xlsx'%filename, index = False)
	aa_for_N_sequon = ['A','R', 'N', 'D', 'B', 'C', 'E', 'Q', 'Z', 'G', 'H', 'I', 'L', 'K', 'M', 'F', 'S', 'T', 'W', 'Y', 'V']
	seq_count = 0
	delete_ind = []
	sequon_lst = []
	double_sequon_seq = []
	for each_pure_seq in pure_seq:
	    each_sequon = []
	    save_num = 0
	    # print('each_pure_seq:%s '%each_pure_seq)
	    for each_aa in aa_for_N_sequon:
	        if 'N%sT'%each_aa in each_pure_seq:
	            for i in range(each_pure_seq.count('N%sT'%each_aa)):
	                each_sequon.append('N%sT'%each_aa)
	                #print('this is each_sequon: %s'%each_sequon)
	                #sequon_lst.append(each_sequon)
	                #print('this is sequon_lst: %s'%sequon_lst)
	                save_num += 1
	        elif 'N%sS'%each_aa in each_pure_seq:
	            for i in range(each_pure_seq.count('N%sS'%each_aa)):
	                each_sequon.append('N%sS'%each_aa)
	                #print('this is each_sequon: %s'%each_sequon)
	                #sequon_lst.append(each_sequon)
	                #print('this is sequon_lst: %s'%sequon_lst)
	                save_num += 1
	    sequon_lst.append(each_sequon)
	    #print('this is sequon_lst: %s'%sequon_lst)
	    # print('save num: %s'%save_num)
	    # print('seq count: %s'%seq_count)
	    #record the key info. from the sequon analysis
	    if save_num == 0: #no sequon
	        delete_ind.append(seq_count)
	    elif save_num == 2: #double sites
	        double_sequon_seq.append(each_pure_seq)
	    seq_count += 1
	    
	# print('\n%s rows will be deleted.'%len(delete_ind))
	sorted_data_scoreHightoLow_score200_pep2d0001 = sorted_data_scoreHightoLow_score200_pep2d0001.drop(sorted_data_scoreHightoLow_score200_pep2d0001.index[delete_ind])
	sorted_data_scoreHightoLow_score200_pep2d0001 = sorted_data_scoreHightoLow_score200_pep2d0001.reset_index(drop=True)
	# print('Data size after deletion:')
	# print(sorted_data_scoreHightoLow_score200_pep2d0001.shape)
	sequon_lst = [sequon[0] if len(sequon) == 1 else sequon for sequon in sequon_lst if len(sequon) != 0]
	#print('this is sequon_lst: %s'%sequon_lst)
	# print('sequon length: %s'%len(sequon_lst))
	sorted_data_scoreHightoLow_score200_pep2d0001.insert(sorted_data_scoreHightoLow_score200_pep2d0001.columns.get_loc('Pure Sequence') + 1 , 'Sequon', sequon_lst , True)

	sorted_data_scoreHightoLow_score200_pep2d0001 = sorted_data_scoreHightoLow_score200_pep2d0001.reset_index(drop=True)

	#store the Pos. numbers & multi-sites positions.
	pos_num = sorted_data_scoreHightoLow_score200_pep2d0001['Pos.'].tolist()

	##sum up the mw. changes for multi-sites.(after dropping rows)
	glycan_pos_num = sorted_data_scoreHightoLow_score200_pep2d0001['Glycans Pos.'].tolist()
	#print('this is glycan_pos_num: %s'%glycan_pos_num)
	multi_ind = []
	glycan_pos_num_count = 0
	for num in glycan_pos_num:
	    #print(num)
	    if type(num) == str:
	        multi_ind.append(glycan_pos_num_count)
	        glycan_pos_num_count += 1
	    else:
	        glycan_pos_num_count += 1
	#print('this is multi_ind: %s'%multi_ind)
	sequence = sorted_data_scoreHightoLow_score200_pep2d0001['Sequence'].tolist() 
	new_seq = []
	for i in range(len(sequence)):
	    if i in multi_ind: #multi-assignment
	        b = re.findall("\d+\.\d+", "Current Level: %s"%sequence[i])
	        b = sum([float(num) for num in b if num != '57.02146' and num != '0.98402' and num != '15.99492'])
	        sequence[i] = sequence[i] + ' --> +%.5f'%(b)
	        new_seq.append(sequence[i])
	    else:
	        new_seq.append(sequence[i])
	new_seq = pd.DataFrame(new_seq)
	sorted_data_scoreHightoLow_score200_pep2d0001['Sequence'] = new_seq
	#display(HTML(sorted_data_scoreHightoLow_score200_pep2d0001.to_html()))

	# print('\nSave the indices of the double-sequon peptide sequence.')
	double_sequon_seq = sorted(set(double_sequon_seq), key=lambda x: double_sequon_seq.index(x))
	#print('this is double_sequon_seq: %s'%double_sequon_seq)
	double_sequon_seq_ind = []
	double_sequon_pos_num = []
	new_glyco_site = []
	for seq in double_sequon_seq:
	    each_double_sequon_pos = []
	    # print('this is seq: %s'%seq)
	    each_double_sequon_seq_ind = sorted_data_scoreHightoLow_score200_pep2d0001.index[sorted_data_scoreHightoLow_score200_pep2d0001['Pure Sequence'] == seq].tolist()
	    #print('this is each_double_sequon_seq_ind: %s'%each_double_sequon_seq_ind)
	    each_sequon = [sequon_lst[i] for i in each_double_sequon_seq_ind][0]
	    # print('this is each_sequon: %s'%each_sequon)
	    if each_sequon[0] == each_sequon[1]:
	        for sequon in each_sequon: 
	            each_N_ind = [m.start() for m in re.finditer(sequon, seq)]       
	    else:
	        each_N_ind = [seq.index(sequon) for sequon in each_sequon]
	        each_N_ind = sorted(each_N_ind)
	        # print('this is each_N_ind: %s'%each_N_ind)
	    each_double_sequon_pos_num = [pos_num[i] for i in each_double_sequon_seq_ind]
	    each_N_ind = np.array(each_N_ind)
	    each_new_glyco_site = [list(each_N_ind+j -1) for j in each_double_sequon_pos_num]
	    new_glyco_site.append(each_new_glyco_site)
	    # print('each_new_glyco_site: %s'%each_new_glyco_site)
	    double_sequon_seq_ind.append(each_double_sequon_seq_ind)
	    double_sequon_pos_num.append(each_double_sequon_pos_num)
	    
	#print('this is double_sequon_seq_ind: %s'%double_sequon_seq_ind)
	#print('this is double_sequon_pos_num: %s'%double_sequon_pos_num)
	#print('new_glyco_site: %s'%new_glyco_site)
	double_sequon_seq_ind = sum(double_sequon_seq_ind, [])
	#print('this is double_sequon_seq_ind: %s'%double_sequon_seq_ind) #for item replacement
	new_glyco_site = sum(new_glyco_site, [])
	#print('this is new_glyco_site: %s'%new_glyco_site)

	##add a column called 'glycosylation site' for real glycan positions. (pos. number + glycan posi)
	glycan_posi_lst = sorted_data_scoreHightoLow_score200_pep2d0001['Glycans Pos.'].tolist()
	glycan_posi_lst = [int(posi) if type(posi) == float else posi for posi in glycan_posi_lst]
	int_glycan_posi_lst = []
	for posi in glycan_posi_lst:
	    #print(posi)
	    #print(type(posi))
	    posi = str(posi)
	    posi = posi.split(';')
	    if len(posi) >1:
	        multiple_site = []
	        for single_posi in posi:
	            single_posi = int(single_posi)
	            multiple_site.append(single_posi)
	        int_glycan_posi_lst.append(multiple_site)
	    elif len(posi) == 1:
	        posi = int(posi[0])
	        int_glycan_posi_lst.append(posi)
	#print('this is int glycan posi lst:\n%s'%int_glycan_posi_lst)

	#add the pos. number to the glycan posi to get the final glycosylation sites. 
	glycosylation_site = []
	int_glycan_posi_count = 0 #make sure the loop went through each item in the int_glycan_posi_lst and keep track the loop. 
	for int_glycan_posi in int_glycan_posi_lst:
	    if type(int_glycan_posi) == list:
	        multiple_glycan_site = []
	        for j in range(len(int_glycan_posi)):
	            each_glycosylation_site = pos_num[int_glycan_posi_count] + (int_glycan_posi[j] -1)
	            multiple_glycan_site.append(each_glycosylation_site)
	        glycosylation_site.append(str(multiple_glycan_site))
	        int_glycan_posi_count += 1
	    elif int_glycan_posi == 0: # for unoccupied sites. (int & list)
	        #print('int_glycan_posi_count: %s'%int_glycan_posi_count)
	        pure_seq = sorted_data_scoreHightoLow_score200_pep2d0001['Pure Sequence'].tolist()
	        seq = pure_seq[int_glycan_posi_count] #find the N site(s) within the peptide seq.
	        #print('seq: %s'%seq)
	        NXS = []
	        NXT = []
	        for var in aa_for_N_sequon:
	            if 'N%sS'%var in seq:
	                NXS = [m.start() for m in re.finditer('N%sS'%var, seq)]
	            elif 'N%sT'%var in seq:
	                NXT = [m.start() for m in re.finditer('N%sT'%var, seq)]
	        if len(NXS) != 0 and len(NXT) != 0:
	            n_posi = sorted((NXS + NXT))
	            unoccupied_multiple_glycan_site = []
	            for j in range(len(n_posi)):
	                each_glycosylation_site = pos_num[int_glycan_posi_count] + (n_posi[j] -1)
	                unoccupied_multiple_glycan_site.append(each_glycosylation_site)
	            glycosylation_site.append(str(unoccupied_multiple_glycan_site))
	            int_glycan_posi_count += 1
	        elif len(NXS) == 0:
	            n_posi = sorted(NXT)
	            if len(n_posi) == 1: #single n site.
	                each_glycosylation_site = pos_num[int_glycan_posi_count] + n_posi[0] -1
	                glycosylation_site.append(each_glycosylation_site)
	                int_glycan_posi_count += 1
	            elif len(n_posi) > 1:
	                unoccupied_multiple_glycan_site = []
	                for j in range(len(n_posi)):
	                    each_glycosylation_site = pos_num[int_glycan_posi_count] + (n_posi[j] -1)
	                    unoccupied_multiple_glycan_site.append(each_glycosylation_site)
	                glycosylation_site.append(str(unoccupied_multiple_glycan_site))
	                int_glycan_posi_count += 1
	        elif len(NXT) == 0:
	            n_posi = sorted(NXS)
	            if len(n_posi) == 1: #single n site.
	                each_glycosylation_site = pos_num[int_glycan_posi_count] + n_posi[0] -1
	                glycosylation_site.append(each_glycosylation_site)
	                int_glycan_posi_count += 1
	            elif len(n_posi) > 1:
	                unoccupied_multiple_glycan_site = []
	                for j in range(len(n_posi)):
	                    each_glycosylation_site = pos_num[int_glycan_posi_count] + (n_posi[j] -1)
	                    unoccupied_multiple_glycan_site.append(each_glycosylation_site)
	                glycosylation_site.append(str(unoccupied_multiple_glycan_site))
	                int_glycan_posi_count += 1
	    elif type(int_glycan_posi) == int and  int_glycan_posi != 0:  
	        each_glycosylation_site = pos_num[int_glycan_posi_count] + int_glycan_posi -1
	        glycosylation_site.append(each_glycosylation_site)
	        int_glycan_posi_count += 1
	regrouped_glyco_site = glycosylation_site.copy()
	# print('this is regrouped_glyco_site: %s'%regrouped_glyco_site)
	# print('this is regrouped_glyco_site size: %s'%len(regrouped_glyco_site))
	new_glyco_site_count = 0
	for i in double_sequon_seq_ind:
	    regrouped_glyco_site[i] = new_glyco_site[new_glyco_site_count]
	    new_glyco_site_count += 1
	# print('this is regrouped_glyco_site: %s'%regrouped_glyco_site)
	# print('this is regrouped_glyco_site size: %s'%len(regrouped_glyco_site))
	##update the glycosylation sites.
	glycosylation_site = [site if type(site) == int else str(site) for site in regrouped_glyco_site]
	# print('this is revised glycosylation_site: %s'%glycosylation_site)
	sorted_data_scoreHightoLow_score200_pep2d0001.insert(2, 'Glycosylation Site', glycosylation_site, True)

	#display(HTML(sorted_data_scoreHightoLow_score200_pep2d0001.to_html()))

	# print('\nStep2: Sort Score within each glycosylation site & get unique Calc.MH while recording PSM.\n')
	site = sorted(set(glycosylation_site), key=lambda x: glycosylation_site.index(x)) 
	# print('this is site ---> %s'%site)
	site_nostr = site.copy()
	each_site_first_lst = []
	str_site = []
	for each_site in site:
	    # print('each_site: %s'%each_site)
	    if type(each_site) != int:
	        str_site.append(each_site)
	        each_site_first = each_site.replace('[', ',').replace(']', ',').split(',')[1] #take the first number to determine the insertion site.
	        each_site_first = int(each_site_first)
	        # print('this is each site first: %s'%each_site_first)
	        each_site_first_lst.append(each_site_first)
	        site_nostr.remove(each_site)
	        # print('site_nostr: %s'%site_nostr)
	site_nostr = list(site_nostr)
	# print('site_nostr: %s'%site_nostr)
	site_nostr.sort()
	# print('site_nostr: %s'%site_nostr)
	site_nostr_new = site_nostr.copy()
	strsite_firstnum_dic = dict(zip(each_site_first_lst, str_site))
	# print('this is the dict before orderedDict: %s'%strsite_firstnum_dic)
	strsite_firstnum_dic = collections.OrderedDict(sorted(strsite_firstnum_dic.items()))
	# print('this is strsite_firstnum_dic: %s'%strsite_firstnum_dic)
	each_site_first_lst = sorted(each_site_first_lst)
	insert_shift = 0
	for i in range(len(each_site_first_lst)):
	    for j in range(len(site_nostr)):
	        # print('this is each_site_first_lst[j]:%s'%each_site_first_lst[i])
	        # print('this is site_nostr[i]:%s'%site_nostr[j])
	        if j != len(site_nostr) -1: #  not the last index.
	            if each_site_first_lst[i] == site_nostr[j]: 
	                site_nostr_new.insert(j +1 +insert_shift , strsite_firstnum_dic.get(each_site_first_lst[i]))
	                # print('site nostr new: %s'%site_nostr_new)
	                insert_shift += 1
	            elif site_nostr[j] < each_site_first_lst[i] < site_nostr[j+1]:
	                site_nostr_new.insert(j +1 +insert_shift , strsite_firstnum_dic.get(each_site_first_lst[i]))
	                # print('site nostr new: %s'%site_nostr_new)
	                insert_shift += 1
	        elif j == len(site_nostr) -1 : #last index            
	            if each_site_first_lst[i] >= site_nostr[j]: 
	                site_nostr_new.append(strsite_firstnum_dic.get(each_site_first_lst[i]))
	                # print('site nostr new: %s'%site_nostr_new)
	                insert_shift += 1
	# print(site_nostr_new)
	frames = []
	for each_site in site_nostr_new:
	    each_sorted_data_scoreHightoLow = sorted_data_scoreHightoLow_score200_pep2d0001[sorted_data_scoreHightoLow_score200_pep2d0001['Glycosylation Site'] == each_site].sort_values('Score', ascending = False)
	    frames.append(each_sorted_data_scoreHightoLow)
	sorted_data_scoreHightoLow_score200_pep2d0001 = pd.concat(frames)
	#print('Export file: ... scoreHightoLow_score200_pep2d0001_allWithSequonAddGlycoSite')
	#sorted_data_scoreHightoLow_score200_pep2d0001.to_excel('%s_sorted_data_scoreHightoLow_score200_pep2d0001_allWithSequonAddGlycoSite.xlsx'%filename, index = False)

	sorted_data_scoreHightoLow_score200_pep2d0001_mhlst = sorted_data_scoreHightoLow_score200_pep2d0001['Calc. MH'].tolist()
	sorted_data_scoreHightoLow_score200_pep2d0001_mh_once = sorted(set(sorted_data_scoreHightoLow_score200_pep2d0001_mhlst), key=lambda x: sorted_data_scoreHightoLow_score200_pep2d0001_mhlst.index(x))
	each_psm_lst = []
	for mh in sorted_data_scoreHightoLow_score200_pep2d0001_mhlst:
	    each_psm = sorted_data_scoreHightoLow_score200_pep2d0001_mhlst.count(mh)
	    each_psm_lst.append(each_psm)
	sorted_data_scoreHightoLow_score200_pep2d0001.insert(len(sorted_data_scoreHightoLow_score200_pep2d0001.columns), "PSM", each_psm_lst, True)
	highest_score_mh_ind = []
	for mh in sorted_data_scoreHightoLow_score200_pep2d0001_mh_once:
	    each_highest_score_mh_ind = sorted_data_scoreHightoLow_score200_pep2d0001.index[sorted_data_scoreHightoLow_score200_pep2d0001['Calc. MH'] == mh].tolist()[0] 
	    highest_score_mh_ind.append(each_highest_score_mh_ind)
	highest_score_mh_ind = []
	for mh in sorted_data_scoreHightoLow_score200_pep2d0001_mh_once:
	    each_highest_score_mh_ind = sorted_data_scoreHightoLow_score200_pep2d0001.index[sorted_data_scoreHightoLow_score200_pep2d0001['Calc. MH'] == mh].tolist()[0] 
	    highest_score_mh_ind.append(each_highest_score_mh_ind)
	sorted_data_scoreHightoLow_score200_pep2d0001_difMH = sorted_data_scoreHightoLow_score200_pep2d0001.loc[highest_score_mh_ind, :]
	#print('Export file: ... scoreHightoLow_score200_pep2d0001_AddPSMdifMH')
	#sorted_data_scoreHightoLow_score200_pep2d0001_difMH.to_excel('%s_scoreHightoLow_score200_pep2d0001_AddPSMdifMH.xlsx'%filename, index = False)

	# print('\nStep3: Start glycan type analysis.')
	glycans = sorted_data_scoreHightoLow_score200_pep2d0001_difMH['Glycans'].tolist()
	types = [glycan.replace(')', ',').replace('(', ',').split(',') if glycan != 0 else glycan for glycan in glycans]

	### FIND THE MULTI-SITES & SINGLE OUT FIRST ###
	sequon = sorted_data_scoreHightoLow_score200_pep2d0001_difMH['Sequon'].tolist()
	#print('this is sequon: %s'%sequon)
	multi_sites_ind = [ i for i in range(len(sequon)) if type(sequon[i]) == list] 
	#print('multi_sites_ind: %s'%multi_sites_ind)

	glycan_data = []
	glycan_data_remain = []
	type_count = 0
	for glycan in types: #within each frag. #form the [a, b, c, d] format first.
	    #print('this is type_count: %s'%type_count)
	    if glycan != 0 and '; HexNAc' not in glycan and type_count not in multi_sites_ind: #real single sites. Enter normal glycan type classification.
	        if 'HexNAc' in glycan:
	            hexnac_num = int(glycan[glycan.index('HexNAc') + 1])
	        else:
	            hexnac_num = 0
	        if 'Hex' in glycan:
	            hex_num = int(glycan[glycan.index('Hex') + 1])
	        else:
	            hex_num = 0
	        if 'Fuc' in glycan:
	            fuc_num = int(glycan[glycan.index('Fuc') + 1])
	        else:
	            fuc_num = 0
	        if 'NeuAc' in glycan:
	            neuac_num = int(glycan[glycan.index('NeuAc') + 1])
	        else:
	            neuac_num = 0
	        each_glycan_data = [hexnac_num, hex_num, fuc_num, neuac_num]
	        each_glycan_data_remain = list(np.array(each_glycan_data) - [2,3,0,0])
	        glycan_data_remain.append(each_glycan_data_remain)
	    elif glycan != 0 and type_count in multi_sites_ind: #[including "fake single sites" assigned by computer] double sites: sum up the glycans from each site, but enter new glycan type classification.
	        if '; HexNAc' not in glycan: #"fake single sites"
	            if 'HexNAc' in glycan:
	                hexnac_num = int(glycan[glycan.index('HexNAc') + 1])
	            else:
	                hexnac_num = 0
	            if 'Hex' in glycan:
	                hex_num = int(glycan[glycan.index('Hex') + 1])
	            else:
	                hex_num = 0
	            if 'Fuc' in glycan:
	                fuc_num = int(glycan[glycan.index('Fuc') + 1])
	            else:
	                fuc_num = 0
	            if 'NeuAc' in glycan:
	                neuac_num = int(glycan[glycan.index('NeuAc') + 1])
	            else:
	                neuac_num = 0
	            all_glycan_num = [hexnac_num, hex_num, fuc_num, neuac_num]
	        else: #double sites also assigned by computer.
	            ##HEX
	            hex_ind = np.array([index for index, element in enumerate(glycan) if element == 'Hex']) #if none, then will return an empty array.
	            hex_num_ind = list(hex_ind + 1) #if none, it will be an empty list. len = 0
	            hex_num = 0
	            for ind in hex_num_ind:
	                hex_num += int(glycan[ind])
	            ##FUC    
	            fuc_ind = np.array([index for index, element in enumerate(glycan) if element == 'Fuc'])
	            fuc_num_ind = list(fuc_ind + 1)
	            fuc_num = 0
	            for ind in fuc_num_ind:
	                fuc_num += int(glycan[ind])
	            ##NEUAC
	            neuac_ind = np.array([index for index, element in enumerate(glycan) if element == 'NeuAc'])
	            neuac_num_ind = list(neuac_ind + 1)
	            neuac_num = 0
	            for ind in neuac_num_ind:
	                neuac_num += int(glycan[ind])
	            ##HEXNAC(at the base of the core, should not be absent)
	            hexnac_num_ind = [glycan.index('HexNAc') + 1, glycan.index('; HexNAc') + 1]
	            hexnac_num = int(glycan[hexnac_num_ind[0]]) + int(glycan[hexnac_num_ind[1]])
	            all_glycan_num = [hexnac_num, hex_num, fuc_num, neuac_num] #after summation of both assignments.
	            
	        multi_site_remain = list(np.array(all_glycan_num) - [4,6,0,0]) 
	        if any(i < 0 for i in multi_site_remain): #check if two cores can form.
	            no_2core_remain = list(np.array(all_glycan_num) - [2,3,0,0])
	            if no_2core_remain[0] == 0 and no_2core_remain[1] == 0:
	                all_glycan_num.append('Multi Site: Only 1 core!')
	                each_glycan_data = all_glycan_num
	            elif any(k < 0 for k in no_2core_remain):
	                all_glycan_num.append('Multi Site: Cannot form any core!')
	                each_glycan_data = all_glycan_num
	            elif no_2core_remain[0] == 0 and no_2core_remain[1] > 0: #only hexose remained.
	                all_glycan_num.append('Multi Site [count PSM once]: Highmannose')
	                each_glycan_data = all_glycan_num
	            elif no_2core_remain[0] > 0: #there are hexnac remained.
	                all_glycan_num.append('Multi Site [count PSM once]: Hybrid/Complex')
	                each_glycan_data = all_glycan_num
	        
	        else: #at least 2 cores.
	            if multi_site_remain[0] == 0 and multi_site_remain[1] == 0:
	                all_glycan_num.append('Multi Site: 2 cores!')
	                each_glycan_data = all_glycan_num
	            elif multi_site_remain[0] == 0 and multi_site_remain[1] > 0: #only hexose remained.
	                all_glycan_num.append('Multi Site [count PSM twice]: Highmannose')
	                each_glycan_data = all_glycan_num
	            elif multi_site_remain[0] > 0: #there are hexnac remained.
	                all_glycan_num.append('Multi Site [count PSM twice]: Hybrid/Complex')
	                each_glycan_data = all_glycan_num
	            
	    elif glycan == 0:
	        each_glycan_data = 0 #unoccupied sites
	        each_glycan_data_remain = 0
	        glycan_data_remain.append(each_glycan_data_remain)
	    
	    type_count += 1
	    glycan_data.append(each_glycan_data)
	    
	#print('this is glycan_data: %s'%glycan_data)
	#print('this is the len of glycan_data: %s'%len(glycan_data))

	#print('this is glycan_data_remain: %s'%glycan_data_remain)  
	#print('this is the len of glycan_data_remain: %s'%len(glycan_data_remain))

	#### SORTING STARTS: the double (multi) sites should be singled out to perform another set of classification method. ####
	to_remain_ind = 0
	glycan_analysis_result = [] #collect all the final analysis results.

	for glycan in glycan_data: #[a, b, c, d], including real singel sites & double sites & unoccupied sites.
	    if glycan == 0: #unoccupied
	        glycan_analysis_result.append(0)
	        to_remain_ind += 1 #to_remian contains unoccupied sites.
	    else: #remember the multi-sites are also here. but now all the glycan type are like [a, b, c, d] since the data from multi-sites will be combined.  
	        if len(glycan) != 5: # skip the double sites & fake single sites since they are already classified. 
	            to_remain_lst = glycan_data_remain[to_remain_ind]
	            negative_count = sum(1 for num in to_remain_lst if num < 0)
	            if negative_count == 0: #all should be >= 0, otherwise the core is not complete. will be separated.
	                ### HIGHMANNOSE & 2400 (Man4) ####
	                if to_remain_lst[0] == 0 and to_remain_lst[1] >= 1 and to_remain_lst[1] <= 6 : 
	                    glycan.append('Man%s'%(glycan[1]))
	                    glycan_analysis_result.append(glycan)
	                    to_remain_ind += 1
	                #### DEAL WITH COMPLEX TYPE FIRST ####
	                ### INTACT LACNAC TYPE: A is dependent on lacnac number & ONLY CORE: 0 HexNac 0 Hex left ####
	                if to_remain_lst[0] == to_remain_lst[1] and to_remain_lst[2] == 0: #no Fuc
	                    if to_remain_lst[0] >= 1 and to_remain_lst[0] <= 4:
	                        if to_remain_lst[0]!= 1:
	                            glycan.append('A%s/A%sB'%(to_remain_lst[0], to_remain_lst[0] -1 ))
	                            glycan_analysis_result.append(glycan)
	                            to_remain_ind += 1
	                        elif to_remain_lst[0] == 1:
	                            glycan.append('A1')
	                            glycan_analysis_result.append(glycan)
	                            to_remain_ind += 1
	                    elif to_remain_lst[0] == 0: # remain 0000/ 0010 / 0011...etc, only the core is present.
	                        glycan.append('Only core!')
	                        glycan_analysis_result.append(glycan)
	                        to_remain_ind += 1
	                    elif to_remain_lst[0] > 4:
	                        glycan.append('lacnac is seldom more than 4!')
	                        glycan.append('A4/A3B')
	                        glycan_analysis_result.append(glycan)
	                        to_remain_ind += 1
	                elif to_remain_lst[0] == to_remain_lst[1] and to_remain_lst[2] != 0: #Fuc is present. can be more than 1.
	                    if to_remain_lst[0] >= 1 and to_remain_lst[0] <= 4:
	                        if to_remain_lst[0]!= 1:
	                            glycan.append('FA%s/FA%sB'%(to_remain_lst[0], to_remain_lst[0] -1 ))
	                            glycan_analysis_result.append(glycan)
	                            to_remain_ind += 1
	                        elif to_remain_lst[0] == 1:
	                            glycan.append('FA1')
	                            glycan_analysis_result.append(glycan)
	                            to_remain_ind += 1
	                    elif to_remain_lst[0] == 0: # remain 0000/ 0010 / 0011...etc, only the core is present.
	                        glycan.append('Only core!')
	                        glycan_analysis_result.append(glycan)
	                        to_remain_ind += 1
	                    elif to_remain_lst[0] > 4:
	                        glycan.append('lacnac is seldom more than 4!')
	                        glycan.append('FA4/A3B')
	                        glycan_analysis_result.append(glycan)
	                        to_remain_ind += 1
	                ### BROKEN LACNAC TYPE: Hex of the remain num is 0 like 10xx, 20xx...etc.mostly be HexNAc. (if it is 0xxx, then it will have lots of Hex, which is HighMannose.)
	                                    # : Also remain 2100, the left out ones are all HexNnac.
	                                    # : Also remain 0100, the left out ones are all Hex.
	                elif to_remain_lst[0] > to_remain_lst[1] and to_remain_lst[2] == 0:#no Fuc
	                    if to_remain_lst[0] >= 1 and to_remain_lst[0] <= 4:
	                        if to_remain_lst[0]!= 1:
	                            glycan.append('A%s/A%sB'%(to_remain_lst[0], to_remain_lst[0] -1 ))
	                            glycan_analysis_result.append(glycan)
	                            to_remain_ind += 1
	                        elif to_remain_lst[0] == 1:
	                            glycan.append('A1')
	                            glycan_analysis_result.append(glycan)
	                            to_remain_ind += 1
	                    else: #more than 4 is still 4 A.
	                        glycan.append('A4/A3B')
	                        glycan_analysis_result.append(glycan)
	                        to_remain_ind += 1
	                elif to_remain_lst[0] > to_remain_lst[1] and to_remain_lst[2] != 0: #Fuc is present.
	                    if to_remain_lst[0] >= 1 and to_remain_lst[0] <= 4:
	                        if to_remain_lst[0]!= 1:
	                            glycan.append('FA%s/FA%sB'%(to_remain_lst[0], to_remain_lst[0] -1 ))
	                            glycan_analysis_result.append(glycan)
	                            to_remain_ind += 1
	                        elif to_remain_lst[0] == 1:
	                            glycan.append('FA1')
	                            glycan_analysis_result.append(glycan)
	                            to_remain_ind += 1
	                    else: #more than 4 is still 4 A.
	                        glycan.append('FA4/FA3B')
	                        glycan_analysis_result.append(glycan)
	                        to_remain_ind += 1
	                #elif to_remain_lst[0] == 0 and to_remain_lst[1] == 1 and to_remain_lst[2] == 0: #no Fuc.
	                #    glycan.append('A1')
	                #    glycan_analysis_result.append(glycan)
	                #    to_remain_ind += 1
	                #elif to_remain_lst[0] == 0 and to_remain_lst[1] == 1 and to_remain_lst[2] != 0: #Fuc is present.
	                #    glycan.append('FA1')
	                #    glycan_analysis_result.append(glycan)
	                #    to_remain_ind += 1
	                elif to_remain_lst[0] == 0 and to_remain_lst[1] >= 7: #e.g. remain 0700 = 2,10,0,0 
	                    glycan.append('potential N-glycan precursor')
	                    glycan_analysis_result.append(glycan)
	                    to_remain_ind += 1
	                    
	                #### ALSO COMPLEX TYPE: more than one lacnac + at least one Hex above the core. e.g. remain 1200 ####
	                elif to_remain_lst[0] < to_remain_lst[1] and to_remain_lst[0] > 2 and to_remain_lst[2] == 0:#no Fuc
	                    if to_remain_lst[0] <= 4:
	                        glycan.append('A%s/A%sB'%(to_remain_lst[0], to_remain_lst[0] -1 ))
	                        glycan_analysis_result.append(glycan)
	                        to_remain_ind += 1
	                    else: #more than 4 is still 4 A.
	                        glycan.append('A4/A3B')
	                        glycan_analysis_result.append(glycan)
	                        to_remain_ind += 1
	                elif to_remain_lst[0] < to_remain_lst[1] and to_remain_lst[0] > 2 and to_remain_lst[2] != 0:#Fuc is present.
	                    if to_remain_lst[0] <= 4:
	                        glycan.append('FA%s/FA%sB'%(to_remain_lst[0], to_remain_lst[0] -1 ))
	                        glycan_analysis_result.append(glycan)
	                        to_remain_ind += 1
	                    else: #more than 4 is still 4 A.
	                        glycan.append('FA4/FA3B')
	                        glycan_analysis_result.append(glycan)
	                        to_remain_ind += 1
	                #### HYBRID TYPE: at least one lacnac + at least one Hex above the core. 4600 up is hybrid ####
	                elif to_remain_lst[0] < to_remain_lst[1] and to_remain_lst[0] >= 1 and to_remain_lst[0] <= 2 and to_remain_lst[2] == 0:#no Fuc
	                    glycan.append('Hybrid')
	                    glycan_analysis_result.append(glycan)
	                    to_remain_ind += 1
	                elif to_remain_lst[0] < to_remain_lst[1] and to_remain_lst[0] >= 1 and to_remain_lst[0] <= 2 and to_remain_lst[2] != 0:#Fuc is present.
	                    glycan.append('Fhybrid')
	                    glycan_analysis_result.append(glycan)
	                    to_remain_ind += 1
	            else:
	                glycan.append('The core is not complete!')
	                glycan_analysis_result.append(glycan)
	                to_remain_ind += 1
	        else: #multi-sites
	            glycan_analysis_result.append(glycan)
	#print('this is glycan_analysis_result: %s'%glycan_analysis_result)
	sorted_data_scoreHightoLow_score200_pep2d0001_difMH.insert(sorted_data_scoreHightoLow_score200_pep2d0001_difMH.columns.get_loc('Glycans') + 1 , 'Glycan Type Analysis Result', glycan_analysis_result , True)
	print('\nExport file: ... scoreHightoLow_score200_pep2d0001_AddPSMdifMHAddGlycanTypeAnalysisColored')

	##color the groups w/ light blue, light green, light yellow.
	def highlight(x):
	    colors = ['#ADD8E6', '#F0FFF0', '#FFFFFF'] #HEX color code.
	    colors_count = 0
	    df = x.copy()
	    site_ind = []
	    for each_site in site_nostr_new:
	        each_site_ind = sorted_data_scoreHightoLow_score200_pep2d0001_difMH.index[sorted_data_scoreHightoLow_score200_pep2d0001_difMH['Glycosylation Site'] == each_site].tolist()
	        site_ind.append(each_site_ind)
	    for ind in site_ind:
	        color_ind = colors_count%3
	        df.loc[ind, :] = 'background-color: %s'%(colors[color_ind])
	        colors_count += 1
	    return df 

	# bytesio to store temporary file-like objects
	buf1 = BytesIO()

	# By setting the 'engine' in the ExcelWriter constructor.
	writer = pd.ExcelWriter(buf1, engine='openpyxl')

	# sorted_data_scoreHightoLow_score200_pep2d0001_difMH.style.apply(highlight, axis=None).to_excel('%s_scoreHightoLow_score200_pep2d0001_AddPSMdifMHAddGlycanTypeAnalysisColored.xlsx'%filename, index = False)
	# sorted_data_scoreHightoLow_score200_pep2d0001_difMH.style.apply(highlight, axis=None).to_excel(writer, index = False)
	sorted_data_scoreHightoLow_score200_pep2d0001_difMH.style.apply(highlight, axis=None).to_excel(writer, index = False)
	writer.save()

	# print('\nFile exported.')

	# print('\nStep4: Summary of the glycan type composition within each site.')
	glycosylation_site = sorted_data_scoreHightoLow_score200_pep2d0001_difMH['Glycosylation Site'].tolist()
	psm = sorted_data_scoreHightoLow_score200_pep2d0001_difMH['PSM'].tolist()
	psm_count = 0

	bar_size = [] #length == site number.
	pie_size = [] #length == site number.
	for site in site_nostr_new: #loop thru non-repeated sites.
	    each_bar_size = [] #length == 16.
	    each_pie_size = [] #length == 4.
	    man9, man8, man7, man6, man5, hybrid, fhybrid, a1, fa1, a2a1b, fa2fa1b, a3a2b, fa3fa2b, a4a3b, fa4fa3b, unoccupied, multi_highman, multi_complex_hybrid = 0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0,0
	    each_site_ind = [pos for pos, char in enumerate(glycosylation_site) if char == site]
	    #print(each_site_ind)
	    #print(glycan_analysis_result[each_site_ind[0]:(each_site_ind[-1] + 1)])
	    each_glycan_type = glycan_analysis_result[each_site_ind[0]:(each_site_ind[-1] + 1)] #all the glycan analysis results within one site.
	    #print(each_glycan_type)
	    for glycan in each_glycan_type:
	        #print('this is glycan: %s'%glycan)
	        if glycan != 0 and len(glycan) == 5:
	            if glycan[4] == 'Man9':
	                man9 += psm[psm_count]
	                psm_count += 1 
	            elif glycan[4] == 'Man8':
	                man8 += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'Man7':
	                man7 += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'Man6':
	                man6 += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'Man5':
	                man5 += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'Man4':
	                psm_count += 1
	            elif glycan[4] == 'Hybrid':
	                hybrid += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'Fhybrid':
	                fhybrid += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'A1':
	                a1 += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'FA1':
	                fa1 += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'A2/A1B':
	                a2a1b += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'FA2/FA1B':
	                fa2fa1b += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'A3/A2B':
	                a3a2b += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'FA3/FA2B':
	                fa3fa2b += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'A4/A3B':
	                a4a3b += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'FA4/FA3B':
	                fa4fa3b += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'The core is not complete!':
	                psm_count += 1
	            elif glycan[4] == 'Only core!':
	                psm_count += 1
	            elif glycan[4] == 'potential N-glycan precursor':
	                psm_count += 1
	            elif glycan[4] == 'Multi Site: Only 1 core!':
	                psm_count += 1
	            elif glycan[4] == 'Multi Site: Cannot form any core!':
	                psm_count += 1
	            elif glycan[4] == 'Multi Site [count PSM once]: Highmannose':
	                multi_highman += psm[psm_count] 
	                psm_count += 1
	            elif glycan[4] == 'Multi Site [count PSM once]: Hybrid/Complex':
	                multi_complex_hybrid += psm[psm_count]
	                psm_count += 1
	            elif glycan[4] == 'Multi Site: 2 cores!':
	                psm_count += 1
	            elif glycan[4] == 'Multi Site [count PSM twice]: Highmannose':
	                multi_highman += psm[psm_count]*2
	                psm_count += 1
	            elif glycan[4] == 'Multi Site [count PSM twice]: Hybrid/Complex':
	                multi_highman += psm[psm_count]*2
	                psm_count += 1
	        elif glycan == 0: #unoccupied
	            unoccupied += psm[psm_count]
	            psm_count += 1
	            
	    if type(site) == str: #multi-sites
	        each_bar_size.extend([multi_highman, multi_complex_hybrid, unoccupied])
	        bar_size.append(each_bar_size)
	        Highman = each_bar_size[0]
	        Hybrid_Complex = each_bar_size[1]
	        Unoccupied = each_bar_size[-1]
	        each_pie_size.extend([Highman, Hybrid_Complex, Unoccupied])
	        pie_size.append(each_pie_size)
	    else: #real single sites.
	        each_bar_size.extend([man9, man8, man7, man6, man5, hybrid, fhybrid, a1, fa1, a2a1b, fa2fa1b, a3a2b, fa3fa2b, a4a3b, fa4fa3b, unoccupied])
	        bar_size.append(each_bar_size)
	        Highman = sum(each_bar_size[:5])
	        Hybrid = sum(each_bar_size[5:7])
	        Complex = sum(each_bar_size[7:15])
	        Unoccupied = each_bar_size[-1]
	        each_pie_size.extend([Highman, Hybrid, Complex, Unoccupied])
	        pie_size.append(each_pie_size)
	# print('this is bar_size:\n%s'%bar_size)
	original_bar_size = bar_size #save for later manipulation. 
	# print('this is pie_size:\n%s'%pie_size)
	original_pie_size = pie_size #save for later manipulation.
	# print('length of the bar_size (also the final site number): %s'%len(bar_size))
	# print('length of the pie_size (also the final site number): %s'%len(pie_size))

	##record glycosylation site indices.
	sites_for_plot = []
	for site in site_nostr_new:
	    if type(site) == int:
	        sites_for_plot.append('N%s'%site)
	    elif type(site) == str and site.count(',') == 1:
	        sites_for_plot.append('N%s/\nN%s'%((site.replace('[', ',').replace(']', ',').replace(' ', ',').split(','))[1], (site.replace('[', ',').replace(']', ',').replace(' ', ',').split(','))[3]))
	    elif type(site) == str and site.count(',') == 2:
	        sites_for_plot.append('N%s/\nN%s/\nN%s'%((site.replace('[', ',').replace(']', ',').replace(' ', ',').split(','))[1], (site.replace('[', ',').replace(']', ',').replace(' ', ',').split(','))[3], (site.replace('[', ',').replace(']', ',').replace(' ', ',').split(','))[5]))
	pie_types = ['HighMannose', 'Hybrid', 'Complex', 'Hybrid/Complex', 'Unoccupied']
	bar_types = ['Man9', 'Man8', 'Man7', 'Man6', 'Man5', 'Hybrid', 'Fhybrid', 'A1', 'FA1', 'A2/A1B', 'FA2/FA1B', 'A3/A2B', 'FA3/FA2B', 'A4/A3B', 'FA4/FA3B', 'Unoccupied']
	multi_bar_types = ['HighMannose', 'Hybrid/Complex', 'Unoccupied']
	colors = ['#008000', '#ffffff','#ff69b4','#d3d3d3'] 
	multi_colors = ['#008000', '#ffc0cb', '#d3d3d3']
	leg_colors = ['#008000', '#ffffff','#ff69b4', '#ffc0cb', '#d3d3d3']

	##start plotting.
	## PLOT BAR & PIE COMBINED CHARTS FIRST.
	##### adjustable parameters ###################
	if len(bar_size) <= 4:
	    number_of_sites_inrow = 1
	    full_row_num = len(bar_size)//number_of_sites_inrow
	    remainder_num = len(bar_size)%number_of_sites_inrow
	    bar_remainder = bar_size[-remainder_num:] #last ones.
	    # print('bar_remainder: %s'%bar_remainder)
	    bar_size = [bar_size[number_of_sites_inrow*i:number_of_sites_inrow*(i+1)] for i in range(full_row_num)]
	    pie_remainder = pie_size[-remainder_num:] #last ones.
	    # print('pie_remainder: %s'%pie_remainder)
	    pie_size = [pie_size[number_of_sites_inrow*i:number_of_sites_inrow*(i+1)] for i in range(full_row_num)]
	    # print('this is resized bar_size: %s'%bar_size)
	    # print('this is resized pie_size: %s'%pie_size)
	else:
	    number_of_sites_inrow = 4
	    full_row_num = len(bar_size)//number_of_sites_inrow
	    remainder_num = len(bar_size)%number_of_sites_inrow
	    bar_remainder = bar_size[-remainder_num:] #last ones.
	    # print('bar_remainder: %s'%bar_remainder)
	    bar_size = [bar_size[number_of_sites_inrow*i:number_of_sites_inrow*(i+1)] for i in range(full_row_num)]
	    bar_size.append(bar_remainder)
	    pie_remainder = pie_size[-remainder_num:] #last ones.
	    # print('pie_remainder: %s'%pie_remainder)
	    pie_size = [pie_size[number_of_sites_inrow*i:number_of_sites_inrow*(i+1)] for i in range(full_row_num)]
	    pie_size.append(pie_remainder)
	    # print('this is resized bar_size: %s'%bar_size)
	    # print('this is resized pie_size: %s'%pie_size)

	if remainder_num == 0:
	    row_num = full_row_num
	else: 
	    row_num = full_row_num + 1

	# print('this is row_num: %s'%row_num)
	# print('this is remainder_num: %s'%remainder_num)
	col_num = 2*number_of_sites_inrow
	fig_ratio = number_of_sites_inrow/row_num
	###############################################

	#make everything in the plot bold & basic setups.
	####### fontsize & spine width #######
	bar_title_fontsize = 12
	SUBPLOT_PAD = 2
	bar_ylab_fontsize = bar_title_fontsize - 4
	bar_width = 0.7
	multi_bar_width = 0.1
	x_tick_fontsize = bar_title_fontsize -4 
	y_tick_fonsize = bar_title_fontsize -4
	x_tick_w = 1.2
	x_tick_l = x_tick_w*3
	y_tick_w = 1.2
	y_tick_l = y_tick_w*3
	spine_width = 1.2
	pie_line_w = 1.2
	bar_line_w = pie_line_w
	legend_fontsize = bar_title_fontsize
	legend_handleL = 2
	MULTI_SITE_BAR_FONT_SIZE = 8
	MULTI_SITE_BAR_LABEL_PAD = 0
	PSM_LABEL_PAD = 1.5
	PIE_R = 1.75
	HATCH = '////'
	TICKER_THRESHOLD = 6
	MAXNLOCATOR_NBINS = 5
	COMBINED_PLOTS_RATIO_1 = 5
	COMBINED_PLOTS_RATIO_2 = 1.5
	######################################
	plt.rcParams["font.weight"] = "bold"
	plt.rcParams["axes.labelweight"] = "bold"
	fig, axs = plt.subplots(nrows=row_num, ncols=col_num, figsize=(COMBINED_PLOTS_RATIO_1*number_of_sites_inrow, COMBINED_PLOTS_RATIO_2*row_num), constrained_layout=True)
	site_name_count = 0

	if row_num >1: 
	    for i in range(row_num):
	        if i != (row_num -1): #other than the last row. 
	            for j in range(number_of_sites_inrow): #N sites from left to right.
	                if max(bar_size[i][j]) > TICKER_THRESHOLD: #use ticker.
	                    if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                        #pie
	                        axs[i, 2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                        axs[i, 2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})[0][1].set_hatch(HATCH)
	                        #bar
	                        axs[i, 2*j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                        axs[i, 2*j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                        yticks = ticker.MaxNLocator(MAXNLOCATOR_NBINS, integer = True)
	                        axs[i, 2*j].yaxis.set_major_locator(yticks)
	                        for axis in ["left", 'bottom']:
	                            axs[i, 2*j].spines[axis].set_linewidth(spine_width)
	                        for spine in ["top", "right"]:
	                            axs[i, 2*j].spines[spine].set_visible(False)
	                        index = [i for i in range(3)]
	                        xticks = np.arange(3)
	                        #yticks = np.arange(max(size[i])+1)
	                        axs[i, 2*j].bar(index[0], bar_size[i][j][0], multi_bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                        axs[i, 2*j].bar(index[1], bar_size[i][j][1], multi_bar_width, color=['#ffc0cb'], edgecolor='black', linewidth = bar_line_w, hatch=HATCH)
	                        axs[i, 2*j].bar(index[2], bar_size[i][j][2], multi_bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                        axs[i, 2*j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                        axs[i, 2*j].title.set_position([-0.15, 1.1])
	                        axs[i, 2*j].set_xticks(xticks, minor=False)
	                        axs[i, 2*j].set_xticklabels('')
	                        #axs[i, 2*j].set_xticklabels(multi_bar_types , fontweight="bold" ,minor=False, rotation = 45, ha = 'center', fontsize = x_tick_fontsize)
	                        #ax1.set_yticks(yticks, minor=False)
	                        #ax1.set_ylim(yticks[0], yticks[-1])
	                        axs[i, 2*j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                        bars = [rect for rect in axs[i, 2*j].get_children() if isinstance(rect, mpl.patches.Rectangle)]
	                        # print('this is the length of bars: %s'%len(bars))
	                        bar_type_count = 0
	                        multi_bar_labels = ['HM', 'H/C', 'UN']
	                        for bar in bars:
	                            if bar_type_count != 3: 
	                                height = bar.get_height() + MULTI_SITE_BAR_LABEL_PAD
	                                label_x_pos = bar.get_x() + bar.get_width() / 2
	                                axs[i, 2*j].text(label_x_pos, height, s=multi_bar_labels[bar_type_count], ha='center',
	                                va='bottom', fontsize = MULTI_SITE_BAR_FONT_SIZE)
	                                bar_type_count += 1
	                            else:
	                                pass
	                        site_name_count += 1
	                    else: #real singel sites.
	                        #pie
	                        axs[i, 2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                        #bar
	                        axs[i, 2*j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                        axs[i, 2*j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                        yticks = ticker.MaxNLocator(MAXNLOCATOR_NBINS, integer = True)
	                        axs[i, 2*j].yaxis.set_major_locator(yticks)
	                        for axis in ["left", 'bottom']:
	                            axs[i, 2*j].spines[axis].set_linewidth(spine_width)
	                        for spine in ["top", "right"]:
	                            axs[i, 2*j].spines[spine].set_visible(False)
	                        index = [i for i in range(16)]
	                        xticks = np.arange(16)
	                        #yticks = np.arange(max(size[i])+1)
	                        axs[i, 2*j].bar(index[0:5], bar_size[i][j][0:5], bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                        axs[i, 2*j].bar(index[5:7], bar_size[i][j][5:7], bar_width, color=['#ffffff'], edgecolor='black', linewidth = bar_line_w)
	                        axs[i, 2*j].bar(index[7:15], bar_size[i][j][7:15], bar_width, color=['#ff69b4'], edgecolor='black', linewidth = bar_line_w)
	                        axs[i, 2*j].bar(index[15], bar_size[i][j][15], bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                        axs[i, 2*j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                        axs[i, 2*j].title.set_position([-0.15, 1.1])
	                        axs[i, 2*j].set_xticks(xticks, minor=False)
	                        axs[i, 2*j].set_xticklabels('')
	                        #ax1.set_yticks(yticks, minor=False)
	                        #ax1.set_ylim(yticks[0], yticks[-1])
	                        axs[i, 2*j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                        site_name_count += 1
	                else: #do not use ticker.
	                    if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                        #pie
	                        axs[i, 2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                        axs[i, 2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})[0][1].set_hatch(HATCH)
	                        #bar
	                        axs[i, 2*j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                        axs[i, 2*j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                        #yticks = ticker.MaxNLocator(6)
	                        #axs[i, 2*j].yaxis.set_major_locator(yticks)
	                        for axis in ["left", 'bottom']:
	                            axs[i, 2*j].spines[axis].set_linewidth(spine_width)
	                        for spine in ["top", "right"]:
	                            axs[i, 2*j].spines[spine].set_visible(False)
	                        index = [i for i in range(3)]
	                        xticks = np.arange(3)
	                        yticks = np.arange(max(bar_size[i][j])+1)
	                        axs[i, 2*j].bar(index[0], bar_size[i][j][0], multi_bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                        axs[i, 2*j].bar(index[1], bar_size[i][j][1], multi_bar_width, color=['#ffc0cb'], edgecolor='black', linewidth = bar_line_w, hatch=HATCH)
	                        axs[i, 2*j].bar(index[2], bar_size[i][j][2], multi_bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                        axs[i, 2*j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                        axs[i, 2*j].title.set_position([-0.15, 1.1])
	                        axs[i, 2*j].set_xticks(xticks, minor=False)
	                        axs[i, 2*j].set_xticklabels('')
	                        #axs[i, 2*j].set_xticklabels(multi_bar_types , fontweight="bold" ,minor=False, rotation = 45, ha = 'center', fontsize = x_tick_fontsize)
	                        axs[i, 2*j].set_yticks(yticks, minor=False)
	                        axs[i, 2*j].set_ylim(yticks[0], yticks[-1])
	                        axs[i, 2*j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                        bars = [rect for rect in axs[i, 2*j].get_children() if isinstance(rect, mpl.patches.Rectangle)]
	                        # print('this is the length of bars: %s'%len(bars))
	                        bar_type_count = 0
	                        multi_bar_labels = ['HM', 'H/C', 'UN']
	                        for bar in bars:
	                            if bar_type_count != 3: 
	                                height = bar.get_height() + MULTI_SITE_BAR_LABEL_PAD 
	                                label_x_pos = bar.get_x() + bar.get_width() / 2
	                                axs[i, 2*j].text(label_x_pos, height, s=multi_bar_labels[bar_type_count], ha='center',
	                                va='bottom', fontsize = MULTI_SITE_BAR_FONT_SIZE)
	                                bar_type_count += 1
	                            else:
	                                pass
	                        site_name_count += 1
	                    else: #real singel sites.
	                        #pie
	                        axs[i, 2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                        #bar
	                        axs[i, 2*j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                        axs[i, 2*j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                        #yticks = ticker.MaxNLocator(6)
	                        #axs[i, 2*j].yaxis.set_major_locator(yticks)
	                        for axis in ["left", 'bottom']:
	                            axs[i, 2*j].spines[axis].set_linewidth(spine_width)
	                        for spine in ["top", "right"]:
	                            axs[i, 2*j].spines[spine].set_visible(False)
	                        index = [i for i in range(16)]
	                        xticks = np.arange(16)
	                        yticks = np.arange(max(bar_size[i][j])+1)
	                        axs[i, 2*j].bar(index[0:5], bar_size[i][j][0:5], bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                        axs[i, 2*j].bar(index[5:7], bar_size[i][j][5:7], bar_width, color=['#ffffff'], edgecolor='black', linewidth = bar_line_w)
	                        axs[i, 2*j].bar(index[7:15], bar_size[i][j][7:15], bar_width, color=['#ff69b4'], edgecolor='black', linewidth = bar_line_w)
	                        axs[i, 2*j].bar(index[15], bar_size[i][j][15], bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                        axs[i, 2*j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                        axs[i, 2*j].title.set_position([-0.15, 1.1])
	                        axs[i, 2*j].set_xticks(xticks, minor=False)
	                        axs[i, 2*j].set_xticklabels('')
	                        axs[i, 2*j].set_yticks(yticks, minor=False)
	                        axs[i, 2*j].set_ylim(yticks[0], yticks[-1])
	                        axs[i, 2*j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                        site_name_count += 1
	        else: #for the last row. Deal with the remainder.
	            if remainder_num == 0: #can be divided by 4.
	                last_row_site_num = number_of_sites_inrow
	            else: #cannot be divided by 4.
	                last_row_site_num = remainder_num
	            # print('this is last_row_site_num: %s'%last_row_site_num)
	            for j in range(number_of_sites_inrow):
	                if j < last_row_site_num:  
	                    # print('this is i,j: %s,%s'%(i,j))
	                    if max(bar_size[i][j]) > TICKER_THRESHOLD: #use ticker.
	                        if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                            #pie
	                            axs[i, 2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                            axs[i, 2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})[0][1].set_hatch(HATCH)
	                            #bar
	                            axs[i, 2*j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                            axs[i, 2*j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                            yticks = ticker.MaxNLocator(MAXNLOCATOR_NBINS, integer = True)
	                            axs[i, 2*j].yaxis.set_major_locator(yticks)
	                            for axis in ["left", 'bottom']:
	                                axs[i, 2*j].spines[axis].set_linewidth(spine_width)
	                            for spine in ["top", "right"]:
	                                axs[i, 2*j].spines[spine].set_visible(False)
	                            index = [i for i in range(3)]
	                            xticks = np.arange(3)
	                            #yticks = np.arange(max(size[i])+1)
	                            axs[i, 2*j].bar(index[0], bar_size[i][j][0], multi_bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, 2*j].bar(index[1], bar_size[i][j][1], multi_bar_width, color=['#ffc0cb'], edgecolor='black', linewidth = bar_line_w, hatch=HATCH)
	                            axs[i, 2*j].bar(index[2], bar_size[i][j][2], multi_bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, 2*j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                            axs[i, 2*j].title.set_position([-0.15, 1.1])
	                            axs[i, 2*j].set_xticks(xticks, minor=False)
	                            axs[i, 2*j].set_xticklabels(multi_bar_types , fontweight="bold" ,minor=False, rotation = 45, ha = 'right', fontsize = x_tick_fontsize)
	                            #ax1.set_yticks(yticks, minor=False)
	                            #ax1.set_ylim(yticks[0], yticks[-1])
	                            axs[i, 2*j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                            bars = [rect for rect in axs[i, 2*j].get_children() if isinstance(rect, mpl.patches.Rectangle)]
	                            # print('this is the length of bars: %s'%len(bars))
	                            bar_type_count = 0
	                            multi_bar_labels = ['HM', 'H/C', 'UN']
	                            for bar in bars:
	                                if bar_type_count != 3: 
	                                    height = bar.get_height() + MULTI_SITE_BAR_LABEL_PAD
	                                    label_x_pos = bar.get_x() + bar.get_width() / 2
	                                    axs[i, 2*j].text(label_x_pos, height, s=multi_bar_labels[bar_type_count], ha='center',
	                                    va='bottom', fontsize = MULTI_SITE_BAR_FONT_SIZE)
	                                    bar_type_count += 1
	                                else:
	                                    pass
	                            site_name_count += 1
	                        else: #real single sites.
	                            #pie
	                            axs[i, 2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                            #bar
	                            axs[i, 2*j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                            axs[i, 2*j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                            yticks = ticker.MaxNLocator(MAXNLOCATOR_NBINS, integer = True)
	                            axs[i, 2*j].yaxis.set_major_locator(yticks)
	                            for axis in ["left", 'bottom']:
	                                axs[i, 2*j].spines[axis].set_linewidth(spine_width)
	                            for spine in ["top", "right"]:
	                                axs[i, 2*j].spines[spine].set_visible(False)
	                            index = [i for i in range(16)]
	                            xticks = np.arange(16)
	                            #yticks = np.arange(max(size[i])+1)
	                            axs[i, 2*j].bar(index[0:5], bar_size[i][j][0:5], bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, 2*j].bar(index[5:7], bar_size[i][j][5:7], bar_width, color=['#ffffff'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, 2*j].bar(index[7:15], bar_size[i][j][7:15], bar_width, color=['#ff69b4'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, 2*j].bar(index[15], bar_size[i][j][15], bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, 2*j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                            axs[i, 2*j].title.set_position([-0.15, 1.1])
	                            axs[i, 2*j].set_xticks(xticks, minor=False)
	                            axs[i, 2*j].set_xticklabels( bar_types, fontweight="bold" ,minor=False, rotation = 90, ha = 'center', fontsize = x_tick_fontsize)
	                            #ax1.set_yticks(yticks, minor=False)
	                            #ax1.set_ylim(yticks[0], yticks[-1])
	                            axs[i, 2*j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                            site_name_count += 1
	                    else: #do not use ticker.
	                        if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                            #pie
	                            axs[i, 2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                            axs[i, 2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})[0][1].set_hatch(HATCH)
	                            #bar
	                            axs[i, 2*j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                            axs[i, 2*j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                            #yticks = ticker.MaxNLocator(6)
	                            #axs[i, 2*j].yaxis.set_major_locator(yticks)
	                            for axis in ["left", 'bottom']:
	                                axs[i, 2*j].spines[axis].set_linewidth(spine_width)
	                            for spine in ["top", "right"]:
	                                axs[i, 2*j].spines[spine].set_visible(False)
	                            index = [i for i in range(3)]
	                            xticks = np.arange(3)
	                            yticks = np.arange(max(bar_size[i][j])+1)
	                            axs[i, 2*j].bar(index[0], bar_size[i][j][0], multi_bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, 2*j].bar(index[1], bar_size[i][j][1], multi_bar_width, color=['#ffc0cb'], edgecolor='black', linewidth = bar_line_w, hatch=HATCH)
	                            axs[i, 2*j].bar(index[2], bar_size[i][j][2], multi_bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, 2*j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                            axs[i, 2*j].title.set_position([-0.15, 1.1])
	                            axs[i, 2*j].set_xticks(xticks, minor=False)
	                            axs[i, 2*j].set_xticklabels(multi_bar_types , fontweight="bold" ,minor=False, rotation = 45, ha = 'right', fontsize = x_tick_fontsize)
	                            axs[i, 2*j].set_yticks(yticks, minor=False)
	                            axs[i, 2*j].set_ylim(yticks[0], yticks[-1])
	                            axs[i, 2*j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                            bars = [rect for rect in axs[i, 2*j].get_children() if isinstance(rect, mpl.patches.Rectangle)]
	                            # print('this is the length of bars: %s'%len(bars))
	                            bar_type_count = 0
	                            multi_bar_labels = ['HM', 'H/C', 'UN']
	                            for bar in bars:
	                                if bar_type_count != 3: 
	                                    height = bar.get_height() + MULTI_SITE_BAR_LABEL_PAD
	                                    label_x_pos = bar.get_x() + bar.get_width() / 2
	                                    axs[i, 2*j].text(label_x_pos, height, s=multi_bar_labels[bar_type_count], ha='center',
	                                    va='bottom', fontsize = MULTI_SITE_BAR_FONT_SIZE)
	                                    bar_type_count += 1
	                                else:
	                                    pass
	                            site_name_count += 1
	                        else: #real single sites.
	                            # print('this is i,j: %s,%s'%(i,j))
	                            #pie
	                            axs[i, 2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                            #bar
	                            axs[i, 2*j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                            axs[i, 2*j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                            #yticks = ticker.MaxNLocator(6)
	                            #axs[i, 0].yaxis.set_major_locator(yticks)
	                            for axis in ["left", 'bottom']:
	                                axs[i, 2*j].spines[axis].set_linewidth(spine_width)
	                            for spine in ["top", "right"]:
	                                axs[i, 2*j].spines[spine].set_visible(False)
	                            index = [i for i in range(16)]
	                            xticks = np.arange(16)
	                            yticks = np.arange(max(bar_size[i][j])+1)
	                            axs[i, 2*j].bar(index[0:5], bar_size[i][j][0:5], bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, 2*j].bar(index[5:7], bar_size[i][j][5:7], bar_width, color=['#ffffff'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, 2*j].bar(index[7:15], bar_size[i][j][7:15], bar_width, color=['#ff69b4'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, 2*j].bar(index[15], bar_size[i][j][15], bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, 2*j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                            axs[i, 2*j].title.set_position([-0.15, 1.1])
	                            axs[i, 2*j].set_xticks(xticks, minor=False)
	                            axs[i, 2*j].set_xticklabels(bar_types , fontweight="bold" ,minor=False, rotation = 90, ha = 'center', fontsize = x_tick_fontsize)
	                            axs[i, 2*j].set_yticks(yticks, minor=False)
	                            axs[i, 2*j].set_ylim(yticks[0], yticks[-1])
	                            axs[i, 2*j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                            site_name_count += 1
	                else: #empty subplot(s). hide spines & ticks & labels.
	                    # print('this is i,j for empty subplots: %s,%s'%(i,j))
	                    axs[i, 2*j].spines["left"].set_visible(False)
	                    axs[i, 2*j].spines["right"].set_visible(False)
	                    axs[i, 2*j].spines["top"].set_visible(False)
	                    axs[i, 2*j].spines["bottom"].set_visible(False)
	                    axs[i, 2*j+1].spines["left"].set_visible(False)
	                    axs[i, 2*j+1].spines["right"].set_visible(False)
	                    axs[i, 2*j+1].spines["top"].set_visible(False)
	                    axs[i, 2*j+1].spines["bottom"].set_visible(False)
	                    axs[i, 2*j].xaxis.set_ticks([])
	                    axs[i, 2*j].yaxis.set_ticks([])
	                    axs[i, 2*j+1].xaxis.set_ticks([])
	                    axs[i, 2*j+1].yaxis.set_ticks([])
	else: #if only one row. (based on the layout setting above, this will happen if there's only one site.)
	    i = 0
	    if remainder_num == 0: #can be divided by 4.
	        last_row_site_num = number_of_sites_inrow
	    else: #cannot be divided by 4.
	        last_row_site_num = remainder_num
	    # print('this is last_row_site_num: %s'%last_row_site_num)
	    for j in range(number_of_sites_inrow):
	        if j < last_row_site_num:  
	            if max(bar_size[i][j]) > TICKER_THRESHOLD: #use ticker.
	                if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                    #pie
	                    axs[2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                    axs[2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})[0][1].set_hatch(HATCH)
	                    #bar
	                    axs[2*j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                    axs[2*j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                    yticks = ticker.MaxNLocator(MAXNLOCATOR_NBINS, integer = True)
	                    axs[2*j].yaxis.set_major_locator(yticks)
	                    for axis in ["left", 'bottom']:
	                        axs[2*j].spines[axis].set_linewidth(spine_width)
	                    for spine in ["top", "right"]:
	                        axs[2*j].spines[spine].set_visible(False)
	                    index = [i for i in range(3)]
	                    xticks = np.arange(3)
	                    #yticks = np.arange(max(size[i])+1)
	                    axs[2*j].bar(index[0], bar_size[i][j][0], multi_bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                    axs[2*j].bar(index[1], bar_size[i][j][1], multi_bar_width, color=['#ffc0cb'], edgecolor='black', linewidth = bar_line_w, hatch=HATCH)
	                    axs[2*j].bar(index[2], bar_size[i][j][2], multi_bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                    axs[2*j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                    axs[2*j].title.set_position([-0.15, 1.1])
	                    axs[2*j].set_xticks(xticks, minor=False)
	                    axs[2*j].set_xticklabels(multi_bar_types , fontweight="bold" ,minor=False, rotation = 45, ha = 'right', fontsize = x_tick_fontsize)
	                    #ax1.set_yticks(yticks, minor=False)
	                    #ax1.set_ylim(yticks[0], yticks[-1])
	                    axs[2*j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                    bars = [rect for rect in axs[2*j].get_children() if isinstance(rect, mpl.patches.Rectangle)]
	                    # print('this is the length of bars: %s'%len(bars))
	                    bar_type_count = 0
	                    multi_bar_labels = ['HM', 'H/C', 'UN']
	                    for bar in bars:
	                        if bar_type_count != 3: 
	                            height = bar.get_height() + MULTI_SITE_BAR_LABEL_PAD
	                            label_x_pos = bar.get_x() + bar.get_width() / 2
	                            axs[2*j].text(label_x_pos, height, s=multi_bar_labels[bar_type_count], ha='center',
	                            va='bottom', fontsize = MULTI_SITE_BAR_FONT_SIZE)
	                            bar_type_count += 1
	                        else:
	                            pass
	                    site_name_count += 1
	                else: #real single sites.
	                    #pie
	                    axs[2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                    #bar
	                    axs[2*j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                    axs[2*j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                    yticks = ticker.MaxNLocator(MAXNLOCATOR_NBINS, integer = True)
	                    axs[2*j].yaxis.set_major_locator(yticks)
	                    for axis in ["left", 'bottom']:
	                        axs[2*j].spines[axis].set_linewidth(spine_width)
	                    for spine in ["top", "right"]:
	                        axs[2*j].spines[spine].set_visible(False)
	                    index = [i for i in range(16)]
	                    xticks = np.arange(16)
	                    #yticks = np.arange(max(size[i])+1)
	                    axs[2*j].bar(index[0:5], bar_size[i][j][0:5], bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                    axs[2*j].bar(index[5:7], bar_size[i][j][5:7], bar_width, color=['#ffffff'], edgecolor='black', linewidth = bar_line_w)
	                    axs[2*j].bar(index[7:15], bar_size[i][j][7:15], bar_width, color=['#ff69b4'], edgecolor='black', linewidth = bar_line_w)
	                    axs[2*j].bar(index[15], bar_size[i][j][15], bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                    axs[2*j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                    axs[2*j].title.set_position([-0.15, 1.1])
	                    axs[2*j].set_xticks(xticks, minor=False)
	                    axs[2*j].set_xticklabels( bar_types, fontweight="bold" ,minor=False, rotation = 90, ha = 'center', fontsize = x_tick_fontsize)
	                    #ax1.set_yticks(yticks, minor=False)
	                    #ax1.set_ylim(yticks[0], yticks[-1])
	                    axs[2*j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                    site_name_count += 1
	            else: #do not use ticker.
	                if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                    #pie
	                    axs[2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                    axs[2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})[0][1].set_hatch(HATCH)
	                    #bar
	                    axs[i, 2*j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                    axs[i, 2*j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                    #yticks = ticker.MaxNLocator(6)
	                    #axs[i, 2*j].yaxis.set_major_locator(yticks)
	                    for axis in ["left", 'bottom']:
	                        axs[2*j].spines[axis].set_linewidth(spine_width)
	                    for spine in ["top", "right"]:
	                        axs[2*j].spines[spine].set_visible(False)
	                    index = [i for i in range(3)]
	                    xticks = np.arange(3)
	                    yticks = np.arange(max(bar_size[i][j])+1)
	                    axs[2*j].bar(index[0], bar_size[i][j][0], multi_bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                    axs[2*j].bar(index[1], bar_size[i][j][1], multi_bar_width, color=['#ffc0cb'], edgecolor='black', linewidth = bar_line_w, hatch=HATCH)
	                    axs[2*j].bar(index[2], bar_size[i][j][2], multi_bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                    axs[2*j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                    axs[2*j].title.set_position([-0.15, 1.1])
	                    axs[2*j].set_xticks(xticks, minor=False)
	                    axs[2*j].set_xticklabels(multi_bar_types , fontweight="bold" ,minor=False, rotation = 45, ha = 'right', fontsize = x_tick_fontsize)
	                    axs[2*j].set_yticks(yticks, minor=False)
	                    axs[2*j].set_ylim(yticks[0], yticks[-1])
	                    axs[2*j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                    bars = [rect for rect in axs[2*j].get_children() if isinstance(rect, mpl.patches.Rectangle)]
	                    # print('this is the length of bars: %s'%len(bars))
	                    bar_type_count = 0
	                    multi_bar_labels = ['HM', 'H/C', 'UN']
	                    for bar in bars:
	                        if bar_type_count != 3: 
	                            height = bar.get_height() + MULTI_SITE_BAR_LABEL_PAD
	                            label_x_pos = bar.get_x() + bar.get_width() / 2
	                            axs[2*j].text(label_x_pos, height, s=multi_bar_labels[bar_type_count], ha='center',
	                            va='bottom', fontsize = MULTI_SITE_BAR_FONT_SIZE)
	                            bar_type_count += 1
	                        else:
	                            pass
	                    site_name_count += 1
	                else: #real single sites.
	                    # print('this is i,j: %s,%s'%(i,j))
	                    #pie
	                    axs[2*j+1].pie(pie_size[i][j], radius=PIE_R, colors = colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                    #bar
	                    axs[2*j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                    axs[2*j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                    #yticks = ticker.MaxNLocator(6)
	                    #axs[i, 0].yaxis.set_major_locator(yticks)
	                    for axis in ["left", 'bottom']:
	                        axs[2*j].spines[axis].set_linewidth(spine_width)
	                    for spine in ["top", "right"]:
	                        axs[2*j].spines[spine].set_visible(False)
	                    index = [i for i in range(16)]
	                    xticks = np.arange(16)
	                    yticks = np.arange(max(bar_size[i][j])+1)
	                    axs[2*j].bar(index[0:5], bar_size[i][j][0:5], bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                    axs[2*j].bar(index[5:7], bar_size[i][j][5:7], bar_width, color=['#ffffff'], edgecolor='black', linewidth = bar_line_w)
	                    axs[2*j].bar(index[7:15], bar_size[i][j][7:15], bar_width, color=['#ff69b4'], edgecolor='black', linewidth = bar_line_w)
	                    axs[2*j].bar(index[15], bar_size[i][j][15], bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                    axs[2*j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                    axs[2*j].title.set_position([-0.15, 1.1])
	                    axs[2*j].set_xticks(xticks, minor=False)
	                    axs[2*j].set_xticklabels(bar_types , fontweight="bold" ,minor=False, rotation = 90, ha = 'center', fontsize = x_tick_fontsize)
	                    axs[2*j].set_yticks(yticks, minor=False)
	                    axs[2*j].set_ylim(yticks[0], yticks[-1])
	                    axs[2*j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                    site_name_count += 1
	        else: #empty subplot(s). hide spines & ticks & labels.
	            # print('this is i,j for empty subplots: %s,%s'%(i,j))
	            axs[2*j].spines["left"].set_visible(False)
	            axs[2*j].spines["right"].set_visible(False)
	            axs[2*j].spines["top"].set_visible(False)
	            axs[2*j].spines["bottom"].set_visible(False)
	            axs[2*j+1].spines["left"].set_visible(False)
	            axs[2*j+1].spines["right"].set_visible(False)
	            axs[2*j+1].spines["top"].set_visible(False)
	            axs[2*j+1].spines["bottom"].set_visible(False)
	            axs[2*j].xaxis.set_ticks([])
	            axs[2*j].yaxis.set_ticks([])
	            axs[2*j+1].xaxis.set_ticks([])
	            axs[2*j+1].yaxis.set_ticks([])


	params = {'legend.fontsize': legend_fontsize, 'legend.handlelength': legend_handleL}
	legend_properties = {'weight':'bold'}
	plt.rcParams.update(params)

	legend_elements = []
	for i in range(len(leg_colors)):
	    if i == 3:
	        each_line = Patch(facecolor = leg_colors[i], edgecolor = 'k', label = pie_types[i], lw = 1.5, hatch = HATCH)
	        legend_elements.append(each_line)
	    else:
	        each_line = Patch(facecolor = leg_colors[i], edgecolor = 'k', label = pie_types[i], lw = 1.5)
	        legend_elements.append(each_line)


	plt.legend(handles=legend_elements, prop=legend_properties, bbox_to_anchor=(1.2,-0.2), loc='upper left')

	#plt.tight_layout()
	# print('\nStep6: Export Pie & Bar charts as .png files.')

	# save plt to bytesio
	in_memory_fp0 = BytesIO()
	plt.savefig(in_memory_fp0)
	plt.close()

	# print('\nStep6: Export All Pie charts as .png')

	# open new ppt object
	prs = Presentation()
	left = top = Inches(1)
	# add 1st ppt slide with pic
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	pic = slide.shapes.add_picture(in_memory_fp0, left, top)
	# fig.savefig('%s_BarPieCharts.png'%filename)

	# store the plot into in-memory
	# fig.savefig(in_memory_fp, format = 'png') 
	# filenames.append('%s_BarPieCharts.png'%filename)

	## DRAW ALL PIE CHARTS (if site num > 4).
	if len(original_pie_size) > 4: 
	##### adjustable parameters ###################
	#if len(original_pie_size) < 4:
	#    number_of_sites_inrow = 1
	#    full_row_num = len(original_pie_size)//number_of_sites_inrow
	#    remainder_num = len(original_pie_size)%number_of_sites_inrow
	#    pie_remainder = original_pie_size[-remainder_num:] #last ones.	
	#    print('pie_remainder: %s'%pie_remainder)
	#    pie_size = [original_pie_size[number_of_sites_inrow*i:number_of_sites_inrow*(i+1)] for i in range(full_row_num)]
	#    print('this is resized pie_size: %s'%pie_size)
	#else:
	    number_of_sites_inrow = 4
	    full_row_num = len(original_pie_size)//number_of_sites_inrow
	    remainder_num = len(original_pie_size)%number_of_sites_inrow
	    pie_remainder = original_pie_size[-remainder_num:] #last ones.
	    # print('pie_remainder: %s'%pie_remainder)
	    pie_size = [original_pie_size[number_of_sites_inrow*i:number_of_sites_inrow*(i+1)] for i in range(full_row_num)]
	    pie_size.append(pie_remainder)
	    # print('this is resized pie_size: %s'%pie_size)

	    if remainder_num == 0:
	        row_num = full_row_num
	    else: #e.g. 3 sites in total.
	        row_num = full_row_num + 1

	    # print('this is row_num: %s'%row_num)
	    # print('this is remainder_num: %s'%remainder_num)
	    col_num = 1*number_of_sites_inrow
	    fig_ratio = number_of_sites_inrow/row_num
	    ###############################################

	    #make everything in the plot bold & basic setups.
	    ####### fontsize & spine width #######
	    pie_title_fontsize = 12
	    SUBPLOT_PAD = 15
	    pie_line_w = 1.2
	    legend_fontsize = bar_title_fontsize
	    legend_handleL = 2
	    PIE_R = 1.75
	    HATCH = '////'
	    PIE_PLOTS_RATIO_1 = 2
	    PIE_PLOTS_RATIO_2 = 2
	    ######################################
	    fig, axs = plt.subplots(nrows=row_num, ncols=col_num, figsize=(PIE_PLOTS_RATIO_1*number_of_sites_inrow, PIE_PLOTS_RATIO_2*row_num), constrained_layout=False)
	    site_name_count = 0

	    if row_num >1: 
	        for i in range(row_num):
	            if i != (row_num -1): #other than the last row. 
	                for j in range(number_of_sites_inrow): #N sites from left to right.
	                    if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                        #pie
	                        axs[i, j].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                        axs[i, j].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})[0][1].set_hatch(HATCH)
	                        axs[i, j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                        site_name_count += 1
	                    else: #real singel sites.
	                        #pie
	                        axs[i, j].pie(pie_size[i][j], radius=PIE_R, colors = colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                        axs[i, j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                        site_name_count += 1
	            else: #for the last row. Deal with the remainder.
	                if remainder_num == 0: #can be divided by 8.
	                    last_row_site_num = number_of_sites_inrow
	                else: #cannot be divided by 8.
	                    last_row_site_num = remainder_num
	                # print('this is last_row_site_num: %s'%last_row_site_num)
	                for j in range(number_of_sites_inrow):
	                    if j < last_row_site_num:  
	                        # print('this is i,j: %s,%s'%(i,j))
	                        if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                            #pie
	                            axs[i, j].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                            axs[i, j].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})[0][1].set_hatch(HATCH)
	                            axs[i, j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                            site_name_count += 1
	                        else: #real single sites.
	                            #pie
	                            axs[i, j].pie(pie_size[i][j], radius=PIE_R, colors = colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                            axs[i, j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                            site_name_count += 1       
	                    else: #empty subplot(s). hide spines & ticks & labels.
	                        # print('this is i,j for empty subplots: %s,%s'%(i,j))
	                        axs[i, j].spines["left"].set_visible(False)
	                        axs[i, j].spines["right"].set_visible(False)
	                        axs[i, j].spines["top"].set_visible(False)
	                        axs[i, j].spines["bottom"].set_visible(False)
	                        axs[i, j].xaxis.set_ticks([])
	                        axs[i, j].yaxis.set_ticks([])
	    else: #if only one row. (based on the layout setting above, this will happen if there's only one site.)
	        i = 0
	        if remainder_num == 0: #can be divided by 8.
	            last_row_site_num = number_of_sites_inrow
	        else: #cannot be divided by 8.
	            last_row_site_num = remainder_num
	        # print('this is last_row_site_num: %s'%last_row_site_num)
	        for j in range(number_of_sites_inrow):
	            if j < last_row_site_num:
	                if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                    #pie
	                    axs[j].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                    axs[j].pie(pie_size[i][j], radius=PIE_R, colors = multi_colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})[0][1].set_hatch(HATCH)
	                    axs[j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                    site_name_count += 1
	                else: #real single sites.
	                    #pie
	                    axs[j].pie(pie_size[i][j], radius=PIE_R, colors = colors, shadow = True, wedgeprops={"edgecolor":"k",'linewidth': pie_line_w, 'antialiased': True})
	                    axs[j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                    site_name_count += 1
	            else: #empty subplot(s). hide spines & ticks & labels.
	                axs[j].spines["left"].set_visible(False)
	                axs[j].spines["right"].set_visible(False)
	                axs[j].spines["top"].set_visible(False)
	                axs[j].spines["bottom"].set_visible(False)
	                axs[j].xaxis.set_ticks([])
	                axs[j].yaxis.set_ticks([])

	    params = {'legend.fontsize': legend_fontsize, 'legend.handlelength': legend_handleL}
	    legend_properties = {'weight':'bold'}
	    plt.rcParams.update(params)

	    legend_elements = []
	    for i in range(len(leg_colors)):
	        if i == 3:
	            each_line = Patch(facecolor = leg_colors[i], edgecolor = 'k', label = pie_types[i], lw = 1.5, hatch = HATCH)
	            legend_elements.append(each_line)
	        else:
	            each_line = Patch(facecolor = leg_colors[i], edgecolor = 'k', label = pie_types[i], lw = 1.5)
	            legend_elements.append(each_line)

	    plt.legend(handles=legend_elements, prop=legend_properties, bbox_to_anchor=(1.2,-0.2), loc='upper left')
	    plt.tight_layout()

	    # save plt to bytesio
	    in_memory_fp1 = BytesIO()
	    plt.savefig(in_memory_fp1)
	    plt.close()

	    # print('\nStep6: Export All Pie charts as .png')

	    # add 2nd ppt slide with pic
	    slide_layout = prs.slide_layouts[1]
	    slide = prs.slides.add_slide(slide_layout)
	    pic = slide.shapes.add_picture(in_memory_fp1, left, top)
	    # prs.save(in_memory_fp, "python_ppt_v1.pptx")

	    # fig.savefig(in_memory_fp, format = 'png')
	    # filenames.append('%s_PieCharts.png'%filename)

	    number_of_sites_inrow = 4
	    full_row_num = len(original_bar_size)//number_of_sites_inrow
	    remainder_num = len(original_bar_size)%number_of_sites_inrow
	    bar_remainder = original_bar_size[-remainder_num:] #last ones.
	    # print('bar_remainder: %s'%bar_remainder)
	    bar_size = [original_bar_size[number_of_sites_inrow*i:number_of_sites_inrow*(i+1)] for i in range(full_row_num)]
	    bar_size.append(bar_remainder)
	    # print('this is resized bar_size: %s'%bar_size)

	    if remainder_num == 0:
	        row_num = full_row_num
	    else: #e.g. 3 sites in total.
	        row_num = full_row_num + 1

	    # print('this is row_num: %s'%row_num)
	    # print('this is remainder_num: %s'%remainder_num)
	    col_num = 1*number_of_sites_inrow
	    fig_ratio = number_of_sites_inrow/row_num
	    ###############################################

	    #make everything in the plot bold & basic setups.
	    ####### fontsize & spine width #######
	    bar_title_fontsize = 15
	    SUBPLOT_PAD = 2
	    bar_ylab_fontsize = bar_title_fontsize - 4
	    bar_width = 0.7
	    multi_bar_width = 0.1
	    x_tick_fontsize = bar_title_fontsize -4 
	    y_tick_fonsize = bar_title_fontsize -4
	    x_tick_w = 1.2
	    x_tick_l = x_tick_w*3
	    y_tick_w = 1.2
	    y_tick_l = y_tick_w*3
	    spine_width = 1.2
	    pie_line_w = 1.2
	    legend_fontsize = bar_title_fontsize
	    legend_handleL = 2
	    MULTI_SITE_BAR_FONT_SIZE = 8
	    MULTI_SITE_BAR_LABEL_PAD = 0
	    PSM_LABEL_PAD = 1.5
	    PIE_R = 1.75
	    HATCH = '////'
	    TICKER_THRESHOLD = 6
	    MAXNLOCATOR_NBINS = 5
	    BAR_PLOTS_RATIO_1 = 5
	    BAR_PLOTS_RATIO_2 = 2.5
	    ######################################
	    plt.rcParams["font.weight"] = "bold"
	    plt.rcParams["axes.labelweight"] = "bold"
	    fig, axs = plt.subplots(nrows=row_num, ncols=col_num, figsize=(BAR_PLOTS_RATIO_1*number_of_sites_inrow, BAR_PLOTS_RATIO_2*row_num), constrained_layout=False)
	    site_name_count = 0

	    if row_num >1: 
	        for i in range(row_num):
	            if i != (row_num -1): #other than the last row. 
	                for j in range(number_of_sites_inrow): #N sites from left to right.
	                    if max(bar_size[i][j]) > TICKER_THRESHOLD: #use ticker.
	                        if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                            #bar
	                            axs[i, j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                            axs[i, j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                            yticks = ticker.MaxNLocator(MAXNLOCATOR_NBINS, integer = True)
	                            axs[i, j].yaxis.set_major_locator(yticks)
	                            for axis in ["left", 'bottom']:
	                                axs[i, j].spines[axis].set_linewidth(spine_width)
	                            for spine in ["top", "right"]:
	                                axs[i, j].spines[spine].set_visible(False)
	                            index = [i for i in range(3)]
	                            xticks = np.arange(3)
	                            #yticks = np.arange(max(size[i])+1)
	                            axs[i, j].bar(index[0], bar_size[i][j][0], multi_bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, j].bar(index[1], bar_size[i][j][1], multi_bar_width, color=['#ffc0cb'], edgecolor='black', linewidth = bar_line_w, hatch=HATCH)
	                            axs[i, j].bar(index[2], bar_size[i][j][2], multi_bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                            axs[i, j].title.set_position([-0.15, 1.1])
	                            axs[i, j].set_xticks(xticks, minor=False)
	                            axs[i, j].set_xticklabels('')
	                            #axs[i, 2*j].set_xticklabels(multi_bar_types , fontweight="bold" ,minor=False, rotation = 45, ha = 'center', fontsize = x_tick_fontsize)
	                            #ax1.set_yticks(yticks, minor=False)
	                            #ax1.set_ylim(yticks[0], yticks[-1])
	                            axs[i, j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                            bars = [rect for rect in axs[i, j].get_children() if isinstance(rect, mpl.patches.Rectangle)]
	                            # print('this is the length of bars: %s'%len(bars))
	                            bar_type_count = 0
	                            multi_bar_labels = ['HM', 'H/C', 'UN']
	                            for bar in bars:
	                                if bar_type_count != 3: 
	                                    height = bar.get_height() + MULTI_SITE_BAR_LABEL_PAD
	                                    label_x_pos = bar.get_x() + bar.get_width() / 2
	                                    axs[i, j].text(label_x_pos, height, s=multi_bar_labels[bar_type_count], ha='center',
	                                    va='bottom', fontsize = MULTI_SITE_BAR_FONT_SIZE)
	                                    bar_type_count += 1
	                                else:
	                                    pass
	                            site_name_count += 1
	                        else: #real singel sites.
	                            #bar
	                            axs[i, j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                            axs[i, j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                            yticks = ticker.MaxNLocator(MAXNLOCATOR_NBINS, integer = True)
	                            axs[i, j].yaxis.set_major_locator(yticks)
	                            for axis in ["left", 'bottom']:
	                                axs[i, j].spines[axis].set_linewidth(spine_width)
	                            for spine in ["top", "right"]:
	                                axs[i, j].spines[spine].set_visible(False)
	                            index = [i for i in range(16)]
	                            xticks = np.arange(16)
	                            #yticks = np.arange(max(size[i])+1)
	                            axs[i, j].bar(index[0:5], bar_size[i][j][0:5], bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, j].bar(index[5:7], bar_size[i][j][5:7], bar_width, color=['#ffffff'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, j].bar(index[7:15], bar_size[i][j][7:15], bar_width, color=['#ff69b4'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, j].bar(index[15], bar_size[i][j][15], bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                            axs[i, j].title.set_position([-0.15, 1.1])
	                            axs[i, j].set_xticks(xticks, minor=False)
	                            axs[i, j].set_xticklabels('')
	                            #ax1.set_yticks(yticks, minor=False)
	                            #ax1.set_ylim(yticks[0], yticks[-1])
	                            axs[i, j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                            site_name_count += 1
	                    else: #do not use ticker.
	                        if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                            #bar
	                            axs[i, j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                            axs[i, j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                            #yticks = ticker.MaxNLocator(6)
	                            #axs[i, 2*j].yaxis.set_major_locator(yticks)
	                            for axis in ["left", 'bottom']:
	                                axs[i, j].spines[axis].set_linewidth(spine_width)
	                            for spine in ["top", "right"]:
	                                axs[i, j].spines[spine].set_visible(False)
	                            index = [i for i in range(3)]
	                            xticks = np.arange(3)
	                            yticks = np.arange(max(bar_size[i][j])+1)
	                            axs[i, j].bar(index[0], bar_size[i][j][0], multi_bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, j].bar(index[1], bar_size[i][j][1], multi_bar_width, color=['#ffc0cb'], edgecolor='black', linewidth = bar_line_w, hatch=HATCH)
	                            axs[i, j].bar(index[2], bar_size[i][j][2], multi_bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                            axs[i, j].title.set_position([-0.15, 1.1])
	                            axs[i, j].set_xticks(xticks, minor=False)
	                            axs[i, j].set_xticklabels('')
	                            #axs[i, 2*j].set_xticklabels(multi_bar_types , fontweight="bold" ,minor=False, rotation = 45, ha = 'center', fontsize = x_tick_fontsize)
	                            axs[i, j].set_yticks(yticks, minor=False)
	                            axs[i, j].set_ylim(yticks[0], yticks[-1])
	                            axs[i, j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                            bars = [rect for rect in axs[i, j].get_children() if isinstance(rect, mpl.patches.Rectangle)]
	                            # print('this is the length of bars: %s'%len(bars))
	                            bar_type_count = 0
	                            multi_bar_labels = ['HM', 'H/C', 'UN']
	                            for bar in bars:
	                                if bar_type_count != 3: 
	                                    height = bar.get_height() + MULTI_SITE_BAR_LABEL_PAD 
	                                    label_x_pos = bar.get_x() + bar.get_width() / 2
	                                    axs[i, j].text(label_x_pos, height, s=multi_bar_labels[bar_type_count], ha='center',
	                                    va='bottom', fontsize = MULTI_SITE_BAR_FONT_SIZE)
	                                    bar_type_count += 1
	                                else:
	                                    pass
	                            site_name_count += 1
	                        else: #real singel sites.
	                            #bar
	                            axs[i, j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                            axs[i, j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                            #yticks = ticker.MaxNLocator(6)
	                            #axs[i, 2*j].yaxis.set_major_locator(yticks)
	                            for axis in ["left", 'bottom']:
	                                axs[i, j].spines[axis].set_linewidth(spine_width)
	                            for spine in ["top", "right"]:
	                                axs[i, j].spines[spine].set_visible(False)
	                            index = [i for i in range(16)]
	                            xticks = np.arange(16)
	                            yticks = np.arange(max(bar_size[i][j])+1)
	                            axs[i, j].bar(index[0:5], bar_size[i][j][0:5], bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, j].bar(index[5:7], bar_size[i][j][5:7], bar_width, color=['#ffffff'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, j].bar(index[7:15], bar_size[i][j][7:15], bar_width, color=['#ff69b4'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, j].bar(index[15], bar_size[i][j][15], bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                            axs[i, j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                            axs[i, j].title.set_position([-0.15, 1.1])
	                            axs[i, j].set_xticks(xticks, minor=False)
	                            axs[i, j].set_xticklabels('')
	                            axs[i, j].set_yticks(yticks, minor=False)
	                            axs[i, j].set_ylim(yticks[0], yticks[-1])
	                            axs[i, j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                            site_name_count += 1
	            else: #for the last row. Deal with the remainder.
	                if remainder_num == 0: #can be divided by 8.
	                    last_row_site_num = number_of_sites_inrow
	                else: #cannot be divided by 8.
	                    last_row_site_num = remainder_num
	                # print('this is last_row_site_num: %s'%last_row_site_num)
	                for j in range(number_of_sites_inrow):
	                    if j < last_row_site_num:  
	                        # print('this is i,j: %s,%s'%(i,j))
	                        if max(bar_size[i][j]) > TICKER_THRESHOLD: #use ticker.
	                            if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                                #bar
	                                axs[i, j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                                axs[i, j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                                yticks = ticker.MaxNLocator(MAXNLOCATOR_NBINS, integer = True)
	                                axs[i, j].yaxis.set_major_locator(yticks)
	                                for axis in ["left", 'bottom']:
	                                    axs[i, j].spines[axis].set_linewidth(spine_width)
	                                for spine in ["top", "right"]:
	                                    axs[i, j].spines[spine].set_visible(False)
	                                index = [i for i in range(3)]
	                                xticks = np.arange(3)
	                                #yticks = np.arange(max(size[i])+1)
	                                axs[i, j].bar(index[0], bar_size[i][j][0], multi_bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                                axs[i, j].bar(index[1], bar_size[i][j][1], multi_bar_width, color=['#ffc0cb'], edgecolor='black', linewidth = bar_line_w, hatch=HATCH)
	                                axs[i, j].bar(index[2], bar_size[i][j][2], multi_bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                                axs[i, j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                                axs[i, j].title.set_position([-0.15, 1.1])
	                                axs[i, j].set_xticks(xticks, minor=False)
	                                axs[i, j].set_xticklabels(multi_bar_types , fontweight="bold" ,minor=False, rotation = 45, ha = 'right', fontsize = x_tick_fontsize)
	                                #ax1.set_yticks(yticks, minor=False)
	                                #ax1.set_ylim(yticks[0], yticks[-1])
	                                axs[i, j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                                bars = [rect for rect in axs[i, 2*j].get_children() if isinstance(rect, mpl.patches.Rectangle)]
	                                # print('this is the length of bars: %s'%len(bars))
	                                bar_type_count = 0
	                                multi_bar_labels = ['HM', 'H/C', 'UN']
	                                for bar in bars:
	                                    if bar_type_count != 3: 
	                                        height = bar.get_height() + MULTI_SITE_BAR_LABEL_PAD
	                                        label_x_pos = bar.get_x() + bar.get_width() / 2
	                                        axs[i, j].text(label_x_pos, height, s=multi_bar_labels[bar_type_count], ha='center',
	                                        va='bottom', fontsize = MULTI_SITE_BAR_FONT_SIZE)
	                                        bar_type_count += 1
	                                    else:
	                                        pass
	                                site_name_count += 1
	                            else: #real single sites.
	                                #bar
	                                axs[i, j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                                axs[i, j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                                yticks = ticker.MaxNLocator(MAXNLOCATOR_NBINS, integer = True)
	                                axs[i, j].yaxis.set_major_locator(yticks)
	                                for axis in ["left", 'bottom']:
	                                    axs[i, j].spines[axis].set_linewidth(spine_width)
	                                for spine in ["top", "right"]:
	                                    axs[i, j].spines[spine].set_visible(False)
	                                index = [i for i in range(16)]
	                                xticks = np.arange(16)
	                                #yticks = np.arange(max(size[i])+1)
	                                axs[i, j].bar(index[0:5], bar_size[i][j][0:5], bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                                axs[i, j].bar(index[5:7], bar_size[i][j][5:7], bar_width, color=['#ffffff'], edgecolor='black', linewidth = bar_line_w)
	                                axs[i, j].bar(index[7:15], bar_size[i][j][7:15], bar_width, color=['#ff69b4'], edgecolor='black', linewidth = bar_line_w)
	                                axs[i, j].bar(index[15], bar_size[i][j][15], bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                                axs[i, j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                                axs[i, j].title.set_position([-0.15, 1.1])
	                                axs[i, j].set_xticks(xticks, minor=False)
	                                axs[i, j].set_xticklabels( bar_types, fontweight="bold" ,minor=False, rotation = 90, ha = 'center', fontsize = x_tick_fontsize)
	                                #ax1.set_yticks(yticks, minor=False)
	                                #ax1.set_ylim(yticks[0], yticks[-1])
	                                axs[i, j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                                site_name_count += 1
	                        else: #do not use ticker.
	                            if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                                #bar
	                                axs[i, j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                                axs[i, j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                                #yticks = ticker.MaxNLocator(6)
	                                #axs[i, 2*j].yaxis.set_major_locator(yticks)
	                                for axis in ["left", 'bottom']:
	                                    axs[i, j].spines[axis].set_linewidth(spine_width)
	                                for spine in ["top", "right"]:
	                                    axs[i, j].spines[spine].set_visible(False)
	                                index = [i for i in range(3)]
	                                xticks = np.arange(3)
	                                yticks = np.arange(max(bar_size[i][j])+1)
	                                axs[i, j].bar(index[0], bar_size[i][j][0], multi_bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                                axs[i, j].bar(index[1], bar_size[i][j][1], multi_bar_width, color=['#ffc0cb'], edgecolor='black', linewidth = bar_line_w, hatch=HATCH)
	                                axs[i, j].bar(index[2], bar_size[i][j][2], multi_bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                                axs[i, j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                                axs[i, j].title.set_position([-0.15, 1.1])
	                                axs[i, j].set_xticks(xticks, minor=False)
	                                axs[i, j].set_xticklabels(multi_bar_types , fontweight="bold" ,minor=False, rotation = 45, ha = 'right', fontsize = x_tick_fontsize)
	                                axs[i, j].set_yticks(yticks, minor=False)
	                                axs[i, j].set_ylim(yticks[0], yticks[-1])
	                                axs[i, j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                                bars = [rect for rect in axs[i, j].get_children() if isinstance(rect, mpl.patches.Rectangle)]
	                                # print('this is the length of bars: %s'%len(bars))
	                                bar_type_count = 0
	                                multi_bar_labels = ['HM', 'H/C', 'UN']
	                                for bar in bars:
	                                    if bar_type_count != 3: 
	                                        height = bar.get_height() + MULTI_SITE_BAR_LABEL_PAD
	                                        label_x_pos = bar.get_x() + bar.get_width() / 2
	                                        axs[i, j].text(label_x_pos, height, s=multi_bar_labels[bar_type_count], ha='center',
	                                        va='bottom', fontsize = MULTI_SITE_BAR_FONT_SIZE)
	                                        bar_type_count += 1
	                                    else:
	                                        pass
	                                site_name_count += 1
	                            else: #real single sites.
	                                # print('this is i,j: %s,%s'%(i,j))
	                                #bar
	                                axs[i, j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                                axs[i, j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                                #yticks = ticker.MaxNLocator(6)
	                                #axs[i, 0].yaxis.set_major_locator(yticks)
	                                for axis in ["left", 'bottom']:
	                                    axs[i, j].spines[axis].set_linewidth(spine_width)
	                                for spine in ["top", "right"]:
	                                    axs[i, j].spines[spine].set_visible(False)
	                                index = [i for i in range(16)]
	                                xticks = np.arange(16)
	                                yticks = np.arange(max(bar_size[i][j])+1)
	                                axs[i, j].bar(index[0:5], bar_size[i][j][0:5], bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                                axs[i, j].bar(index[5:7], bar_size[i][j][5:7], bar_width, color=['#ffffff'], edgecolor='black', linewidth = bar_line_w)
	                                axs[i, j].bar(index[7:15], bar_size[i][j][7:15], bar_width, color=['#ff69b4'], edgecolor='black', linewidth = bar_line_w)
	                                axs[i, j].bar(index[15], bar_size[i][j][15], bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                                axs[i, j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                                axs[i, j].title.set_position([-0.15, 1.1])
	                                axs[i, j].set_xticks(xticks, minor=False)
	                                axs[i, j].set_xticklabels(bar_types , fontweight="bold" ,minor=False, rotation = 90, ha = 'center', fontsize = x_tick_fontsize)
	                                axs[i, j].set_yticks(yticks, minor=False)
	                                axs[i, j].set_ylim(yticks[0], yticks[-1])
	                                axs[i, j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                                site_name_count += 1
	                    else: #empty subplot(s). hide spines & ticks & labels.
	                        # print('this is i,j for empty subplots: %s,%s'%(i,j))
	                        axs[i, j].spines["left"].set_visible(False)
	                        axs[i, j].spines["right"].set_visible(False)
	                        axs[i, j].spines["top"].set_visible(False)
	                        axs[i, j].spines["bottom"].set_visible(False)
	                        axs[i, j].xaxis.set_ticks([])
	                        axs[i, j].yaxis.set_ticks([])
	    else: #if only one row. (based on the layout setting above, this will happen if there's only one site.)
	        i = 0
	        if remainder_num == 0: #can be divided by 4.
	            last_row_site_num = number_of_sites_inrow
	        else: #cannot be divided by 4.
	            last_row_site_num = remainder_num
	        # print('this is last_row_site_num: %s'%last_row_site_num)
	        for j in range(number_of_sites_inrow):
	            if j < last_row_site_num:  
	                if max(bar_size[i][j]) > TICKER_THRESHOLD: #use ticker.
	                    if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                        #bar
	                        axs[j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                        axs[j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                        yticks = ticker.MaxNLocator(MAXNLOCATOR_NBINS, integer = True)
	                        axs[j].yaxis.set_major_locator(yticks)
	                        for axis in ["left", 'bottom']:
	                            axs[j].spines[axis].set_linewidth(spine_width)
	                        for spine in ["top", "right"]:
	                            axs[j].spines[spine].set_visible(False)
	                        index = [i for i in range(3)]
	                        xticks = np.arange(3)
	                        #yticks = np.arange(max(size[i])+1)
	                        axs[j].bar(index[0], bar_size[i][j][0], multi_bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                        axs[j].bar(index[1], bar_size[i][j][1], multi_bar_width, color=['#ffc0cb'], edgecolor='black', linewidth = bar_line_w, hatch=HATCH)
	                        axs[j].bar(index[2], bar_size[i][j][2], multi_bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                        axs[j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                        axs[j].title.set_position([-0.15, 1.1])
	                        axs[j].set_xticks(xticks, minor=False)
	                        axs[j].set_xticklabels(multi_bar_types , fontweight="bold" ,minor=False, rotation = 45, ha = 'right', fontsize = x_tick_fontsize)
	                        #ax1.set_yticks(yticks, minor=False)
	                        #ax1.set_ylim(yticks[0], yticks[-1])
	                        axs[j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                        bars = [rect for rect in axs[j].get_children() if isinstance(rect, mpl.patches.Rectangle)]
	                        # print('this is the length of bars: %s'%len(bars))
	                        bar_type_count = 0
	                        multi_bar_labels = ['HM', 'H/C', 'UN']
	                        for bar in bars:
	                            if bar_type_count != 3: 
	                                height = bar.get_height() + MULTI_SITE_BAR_LABEL_PAD
	                                label_x_pos = bar.get_x() + bar.get_width() / 2
	                                axs[j].text(label_x_pos, height, s=multi_bar_labels[bar_type_count], ha='center',
	                                va='bottom', fontsize = MULTI_SITE_BAR_FONT_SIZE)
	                                bar_type_count += 1
	                            else:
	                                pass
	                        site_name_count += 1
	                    else: #real single sites.
	                        #bar
	                        axs[j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                        axs[j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                        yticks = ticker.MaxNLocator(MAXNLOCATOR_NBINS, integer = True)
	                        axs[j].yaxis.set_major_locator(yticks)
	                        for axis in ["left", 'bottom']:
	                            axs[j].spines[axis].set_linewidth(spine_width)
	                        for spine in ["top", "right"]:
	                            axs[j].spines[spine].set_visible(False)
	                        index = [i for i in range(16)]
	                        xticks = np.arange(16)
	                        #yticks = np.arange(max(size[i])+1)
	                        axs[j].bar(index[0:5], bar_size[i][j][0:5], bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                        axs[j].bar(index[5:7], bar_size[i][j][5:7], bar_width, color=['#ffffff'], edgecolor='black', linewidth = bar_line_w)
	                        axs[j].bar(index[7:15], bar_size[i][j][7:15], bar_width, color=['#ff69b4'], edgecolor='black', linewidth = bar_line_w)
	                        axs[j].bar(index[15], bar_size[i][j][15], bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                        axs[j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                        axs[j].title.set_position([-0.15, 1.1])
	                        axs[j].set_xticks(xticks, minor=False)
	                        axs[j].set_xticklabels( bar_types, fontweight="bold" ,minor=False, rotation = 90, ha = 'center', fontsize = x_tick_fontsize)
	                        #ax1.set_yticks(yticks, minor=False)
	                        #ax1.set_ylim(yticks[0], yticks[-1])
	                        axs[j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                        site_name_count += 1
	                else: #do not use ticker.
	                    if '/' in sites_for_plot[site_name_count]: #multi-sites.
	                        #bar
	                        axs[j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                        axs[j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                        #yticks = ticker.MaxNLocator(6)
	                        #axs[i, 2*j].yaxis.set_major_locator(yticks)
	                        for axis in ["left", 'bottom']:
	                            axs[j].spines[axis].set_linewidth(spine_width)
	                        for spine in ["top", "right"]:
	                            axs[j].spines[spine].set_visible(False)
	                        index = [i for i in range(3)]
	                        xticks = np.arange(3)
	                        yticks = np.arange(max(bar_size[i][j])+1)
	                        axs[j].bar(index[0], bar_size[i][j][0], multi_bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                        axs[j].bar(index[1], bar_size[i][j][1], multi_bar_width, color=['#ffc0cb'], edgecolor='black', linewidth = bar_line_w, hatch=HATCH)
	                        axs[j].bar(index[2], bar_size[i][j][2], multi_bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                        axs[j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                        axs[j].title.set_position([-0.15, 1.1])
	                        axs[j].set_xticks(xticks, minor=False)
	                        axs[j].set_xticklabels(multi_bar_types , fontweight="bold" ,minor=False, rotation = 45, ha = 'right', fontsize = x_tick_fontsize)
	                        axs[j].set_yticks(yticks, minor=False)
	                        axs[j].set_ylim(yticks[0], yticks[-1])
	                        axs[j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                        bars = [rect for rect in axs[j].get_children() if isinstance(rect, mpl.patches.Rectangle)]
	                        # print('this is the length of bars: %s'%len(bars))
	                        bar_type_count = 0
	                        multi_bar_labels = ['HM', 'H/C', 'UN']
	                        for bar in bars:
	                            if bar_type_count != 3: 
	                                height = bar.get_height() + MULTI_SITE_BAR_LABEL_PAD
	                                label_x_pos = bar.get_x() + bar.get_width() / 2
	                                axs[j].text(label_x_pos, height, s=multi_bar_labels[bar_type_count], ha='center',
	                                va='bottom', fontsize = MULTI_SITE_BAR_FONT_SIZE)
	                                bar_type_count += 1
	                            else:
	                                pass
	                        site_name_count += 1
	                    else: #real single sites.
	                        # print('this is i,j: %s,%s'%(i,j))
	                        #bar
	                        axs[j].xaxis.set_tick_params(width=x_tick_w, length=x_tick_l)
	                        axs[j].yaxis.set_tick_params(width=y_tick_w, length=y_tick_l, labelsize = y_tick_fonsize)
	                        #yticks = ticker.MaxNLocator(6)
	                        #axs[i, 0].yaxis.set_major_locator(yticks)
	                        for axis in ["left", 'bottom']:
	                            axs[j].spines[axis].set_linewidth(spine_width)
	                        for spine in ["top", "right"]:
	                            axs[j].spines[spine].set_visible(False)
	                        index = [i for i in range(16)]
	                        xticks = np.arange(16)
	                        yticks = np.arange(max(bar_size[i][j])+1)
	                        axs[j].bar(index[0:5], bar_size[i][j][0:5], bar_width, color=['#008000'], edgecolor='black', linewidth = bar_line_w)
	                        axs[j].bar(index[5:7], bar_size[i][j][5:7], bar_width, color=['#ffffff'], edgecolor='black', linewidth = bar_line_w)
	                        axs[j].bar(index[7:15], bar_size[i][j][7:15], bar_width, color=['#ff69b4'], edgecolor='black', linewidth = bar_line_w)
	                        axs[j].bar(index[15], bar_size[i][j][15], bar_width, color=['#d3d3d3'], edgecolor='black', linewidth = bar_line_w)
	                        axs[j].set_title(sites_for_plot[site_name_count], size = bar_title_fontsize, fontweight="bold", pad = SUBPLOT_PAD, ha='right')
	                        axs[j].title.set_position([-0.15, 1.1])
	                        axs[j].set_xticks(xticks, minor=False)
	                        axs[j].set_xticklabels(bar_types , fontweight="bold" ,minor=False, rotation = 90, ha = 'center', fontsize = x_tick_fontsize)
	                        axs[j].set_yticks(yticks, minor=False)
	                        axs[j].set_ylim(yticks[0], yticks[-1])
	                        axs[j].set_ylabel("PSM", fontsize = bar_ylab_fontsize, fontweight="bold", labelpad = PSM_LABEL_PAD)
	                        site_name_count += 1
	            else: #empty subplot(s). hide spines & ticks & labels.
	                # print('this is i,j for empty subplots: %s,%s'%(i,j))
	                axs[j].spines["left"].set_visible(False)
	                axs[j].spines["right"].set_visible(False)
	                axs[j].spines["top"].set_visible(False)
	                axs[j].spines["bottom"].set_visible(False)
	                axs[j].xaxis.set_ticks([])
	                axs[j].yaxis.set_ticks([])


	    params = {'legend.fontsize': legend_fontsize, 'legend.handlelength': legend_handleL}
	    legend_properties = {'weight':'bold'}
	    plt.rcParams.update(params)

	    legend_elements = []
	    for i in range(len(leg_colors)):
	        if i == 3:
	            each_line = Patch(facecolor = leg_colors[i], edgecolor = 'k', label = pie_types[i], lw = 1.5, hatch = HATCH)
	            legend_elements.append(each_line)
	        else:
	            each_line = Patch(facecolor = leg_colors[i], edgecolor = 'k', label = pie_types[i], lw = 1.5)
	            legend_elements.append(each_line)


	    plt.legend(handles=legend_elements, prop=legend_properties, bbox_to_anchor=(1.2,-0.2), loc='upper left')

	    plt.tight_layout()

	    # save plt to bytesio
	    in_memory_fp2 = BytesIO()
	    plt.savefig(in_memory_fp2)
	    plt.close()

	    # add 3rd ppt slide with pic
	    slide_layout = prs.slide_layouts[1]
	    slide = prs.slides.add_slide(slide_layout)
	    pic = slide.shapes.add_picture(in_memory_fp2, left, top)

	    # save the ppt object to bytesio
	    buf2 = BytesIO()
	    prs.save(buf2)

	    # print('\nStep6: Export All Bar charts as .png')
	    # fig.savefig(in_memory_fp, format = 'png')

	    # filenames.append('%s_BarCharts.png'%filename)

    # Folder name in ZIP archive 
    zip_filename = "Results.zip"
    # print("Creating archive: {:s}".format(zip_filename))

    zip_buffer = BytesIO()

    with zipfile.ZipFile(zip_buffer, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
    	# 2 in-memory files in total.
    	fname1 = "%s_Colored.xlsx"%filename
    	# print("  Writing {:s} in the archive".format(fname1))
    	zf.writestr(fname1, buf1.getvalue())
    	fname2 = "Plots.pptx"
    	# print("  Writing {:s} in the archive".format(fname2))
    	zf.writestr(fname2, buf2.getvalue())

	return zip_buffer

def get_results(task):
	print(task.result()) 